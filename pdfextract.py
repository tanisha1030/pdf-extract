import fitz  # PyMuPDF
import os
import json
import pandas as pd
from PIL import Image
import pytesseract
from collections import Counter
import zipfile
import tempfile
import warnings
import re
from concurrent.futures import ThreadPoolExecutor, ProcessPoolExecutor
import multiprocessing
from functools import partial
import time
warnings.filterwarnings('ignore')

# Try importing optional libraries with fallbacks
try:
    import pdfplumber
    HAS_PDFPLUMBER = True
except ImportError:
    HAS_PDFPLUMBER = False
    print("Warning: pdfplumber not available. Some table extraction features may be limited.")

try:
    import tabula
    HAS_TABULA = True
except ImportError:
    HAS_TABULA = False
    print("Warning: tabula-py not available. Some table extraction features may be limited.")

try:
    import camelot
    HAS_CAMELOT = True
except ImportError:
    HAS_CAMELOT = False
    print("Warning: camelot-py not available. Some table extraction features may be limited.")

try:
    import cv2
    import numpy as np
    HAS_CV2 = True
except ImportError:
    HAS_CV2 = False
    print("Warning: OpenCV not available. Some image processing features may be limited.")

try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False
    print("Warning: python-docx not available. DOCX processing will be skipped.")

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("Warning: python-pptx not available. PPTX processing will be skipped.")

try:
    import openpyxl
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False
    print("Warning: openpyxl not available. Excel processing may be limited.")

class ComprehensiveDocumentExtractor:
    def __init__(self):
        self.extracted_data = {}
        self.supported_formats = ['.pdf', '.docx', '.pptx', '.xlsx', '.xls']
        self.setup_directories()
        self.max_workers = min(multiprocessing.cpu_count(), 4)  # Limit to avoid memory issues
    
    def setup_directories(self):
        """Create necessary directories"""
        os.makedirs('extracted_images', exist_ok=True)
        os.makedirs('extracted_tables', exist_ok=True)
    
    def extract_pdf_comprehensive(self, pdf_path):
        """Extract comprehensive data from PDF with parallel processing"""
        try:
            print(f"📄 Processing PDF: {os.path.basename(pdf_path)}")
            doc = fitz.open(pdf_path)
            
            pdf_data = {
                'filename': os.path.basename(pdf_path),
                'file_type': 'PDF',
                'metadata': {},
                'pages': [],
                'total_images': 0,
                'total_tables': 0,
                'fonts_used': [],
                'total_words': 0,
                'total_characters': 0,
                'page_count': len(doc),
                'file_size_mb': round(os.path.getsize(pdf_path) / (1024*1024), 2)
            }
            
            # Extract metadata safely
            try:
                metadata = doc.metadata
                if metadata:
                    pdf_data['metadata'] = {k: v for k, v in metadata.items() if v}
            except:
                pdf_data['metadata'] = {}
            
            # Process pages in parallel chunks
            pages_per_chunk = max(1, len(doc) // self.max_workers)
            page_chunks = []
            
            for i in range(0, len(doc), pages_per_chunk):
                chunk_end = min(i + pages_per_chunk, len(doc))
                page_chunks.append((i, chunk_end))
            
            all_pages = []
            all_fonts = []
            
            # Process chunks in parallel
            with ThreadPoolExecutor(max_workers=self.max_workers) as executor:
                chunk_results = list(executor.map(
                    partial(self.process_page_chunk, doc, pdf_path),
                    page_chunks
                ))
            
            # Combine results
            for chunk_pages, chunk_fonts in chunk_results:
                all_pages.extend(chunk_pages)
                all_fonts.extend(chunk_fonts)
            
            pdf_data['pages'] = all_pages
            pdf_data['fonts_used'] = list(set(all_fonts))
            
            # Calculate totals
            for page_data in pdf_data['pages']:
                pdf_data['total_images'] += len(page_data['images'])
                pdf_data['total_tables'] += len(page_data['tables'])
                pdf_data['total_words'] += page_data['word_count']
                pdf_data['total_characters'] += page_data['char_count']
            
            # Extract tables using multiple methods (parallel)
            print("  📊 Extracting tables...")
            pdf_data['extracted_tables'] = self.extract_tables_multiple_methods_parallel(pdf_path)
            
            doc.close()
            print(f"  ✅ PDF processed: {pdf_data['page_count']} pages, {pdf_data['total_words']} words")
            return pdf_data
            
        except Exception as e:
            print(f"❌ Error processing PDF {pdf_path}: {str(e)}")
            return None
    
    def process_page_chunk(self, doc, pdf_path, chunk_range):
        """Process a chunk of pages in parallel"""
        start_page, end_page = chunk_range
        chunk_pages = []
        chunk_fonts = []
        
        for page_num in range(start_page, end_page):
            try:
                page = doc[page_num]
                page_data = self.extract_page_data_optimized(page, page_num + 1, pdf_path)
                chunk_pages.append(page_data)
                chunk_fonts.extend(page_data.get('fonts', []))
                
            except Exception as e:
                print(f"  ⚠️  Error on page {page_num + 1}: {str(e)}")
                continue
        
        return chunk_pages, chunk_fonts
    
    def extract_page_data_optimized(self, page, page_num, pdf_path):
        """Extract data from a single page with optimizations"""
        page_data = {
            'page_number': page_num,
            'text': '',
            'formatted_text': [],
            'images': [],
            'tables': [],
            'fonts': [],
            'word_count': 0,
            'char_count': 0,
            'page_size': {
                'width': round(page.rect.width, 2),
                'height': round(page.rect.height, 2)
            },
            'rotation': page.rotation
        }
        
        try:
            # Extract text with formatting - optimized
            text_dict = page.get_text("dict")
            page_text_parts = []
            fonts_on_page = set()  # Use set for faster lookups
            
            for block in text_dict.get("blocks", []):
                if "lines" in block:
                    for line in block["lines"]:
                        line_text = ""
                        for span in line["spans"]:
                            text = span.get("text", "")
                            line_text += text
                            
                            # Store formatted text with details (reduced data)
                            font_name = span.get('font', 'Unknown')
                            fonts_on_page.add(font_name)
                            
                            # Only store non-empty text spans
                            if text.strip():
                                page_data['formatted_text'].append({
                                    'text': text,
                                    'font': font_name,
                                    'size': round(span.get('size', 0), 1),  # Reduced precision
                                    'bbox': [round(x, 1) for x in span.get('bbox', [0, 0, 0, 0])]
                                })
                        
                        page_text_parts.append(line_text)
            
            page_text = ''.join(page_text_parts)
            page_data['text'] = page_text
            page_data['fonts'] = list(fonts_on_page)
            page_data['word_count'] = len(page_text.split())
            page_data['char_count'] = len(page_text)
            
            # Extract images and tables in parallel
            futures = []
            with ThreadPoolExecutor(max_workers=2) as executor:
                image_future = executor.submit(self.extract_images_from_page_optimized, page, page_num, pdf_path)
                table_future = executor.submit(self.extract_tables_from_page_optimized, page)
                
                page_data['images'] = image_future.result()
                page_data['tables'] = table_future.result()
            
        except Exception as e:
            print(f"    ⚠️  Error extracting from page {page_num}: {str(e)}")
        
        return page_data
    
    def extract_images_from_page_optimized(self, page, page_num, pdf_path):
        """Extract images from a page with optimizations"""
        images = []
        
        try:
            image_list = page.get_images()
            
            if not image_list:
                return images
            
            # Process images in batches
            batch_size = 5
            for i in range(0, len(image_list), batch_size):
                batch = image_list[i:i+batch_size]
                
                for img_index, img in enumerate(batch):
                    try:
                        xref = img[0]
                        base_image = page.parent.extract_image(xref)
                        
                        # Skip very small images (likely artifacts)
                        if base_image["width"] < 50 or base_image["height"] < 50:
                            continue
                        
                        image_bytes = base_image["image"]
                        
                        # Create safe filename
                        base_name = os.path.splitext(os.path.basename(pdf_path))[0]
                        img_name = f"extracted_images/{base_name}_p{page_num}_i{i+img_index+1}.png"
                        
                        # Save image
                        with open(img_name, "wb") as img_file:
                            img_file.write(image_bytes)
                        
                        # Get image info (reduced data)
                        img_info = {
                            'filename': img_name,
                            'width': base_image["width"],
                            'height': base_image["height"],
                            'size_bytes': len(image_bytes)
                        }
                        
                        images.append(img_info)
                        
                    except Exception as e:
                        continue
                        
        except Exception as e:
            pass
        
        return images
    
    def extract_tables_from_page_optimized(self, page):
        """Optimized table detection from page"""
        tables = []
        try:
            text_dict = page.get_text("dict")
            
            # Pre-filter blocks by size and content
            potential_table_blocks = []
            for block in text_dict.get("blocks", []):
                if "lines" in block and len(block["lines"]) >= 3:
                    potential_table_blocks.append(block)
            
            # Process potential table blocks
            for block in potential_table_blocks:
                lines = block["lines"]
                text_lines = []
                
                for line in lines:
                    row_text = ""
                    for span in line["spans"]:
                        row_text += span.get("text", "") + " "
                    text_lines.append(row_text.strip())
                
                # Filter out empty lines
                text_lines = [line for line in text_lines if line.strip()]
                
                # Quick table detection
                if len(text_lines) >= 2 and self.is_likely_table_fast(text_lines):
                    tables.append({
                        'bbox': block.get('bbox', [0, 0, 0, 0]),
                        'data': text_lines,
                        'method': 'fast_detection',
                        'confidence': 'medium'
                    })
            
        except Exception as e:
            pass
        
        return tables
    
    def is_likely_table_fast(self, text_lines):
        """Fast heuristic to determine if text lines represent a table"""
        if len(text_lines) < 2:
            return False
        
        # Quick checks only
        score = 0
        
        # Check for consistent separators
        separator_patterns = [r'\t', r'\s{2,}', r'\|']
        for pattern in separator_patterns:
            counts = [len(re.findall(pattern, line)) for line in text_lines[:5]]  # Check first 5 lines only
            if len(set(counts)) == 1 and counts[0] > 0:
                score += 2
                break
        
        # Check for numeric content
        numeric_lines = sum(1 for line in text_lines[:5] if re.search(r'\d+', line))
        if numeric_lines >= 2:
            score += 1
        
        # Exclude very long lines (likely paragraphs)
        if any(len(line) > 150 for line in text_lines[:3]):
            score -= 2
        
        return score >= 2
    
    def extract_tables_multiple_methods_parallel(self, pdf_path):
        """Extract tables using multiple methods in parallel"""
        all_tables = []
        
        # Run table extraction methods in parallel
        with ThreadPoolExecutor(max_workers=3) as executor:
            futures = []
            
            if HAS_TABULA:
                futures.append(executor.submit(self.extract_tabula_tables, pdf_path))
            
            if HAS_CAMELOT:
                futures.append(executor.submit(self.extract_camelot_tables, pdf_path))
            
            if HAS_PDFPLUMBER:
                futures.append(executor.submit(self.extract_pdfplumber_tables, pdf_path))
            
            # Collect results
            for future in futures:
                try:
                    tables = future.result(timeout=30)  # 30 second timeout per method
                    all_tables.extend(tables)
                except Exception as e:
                    continue
        
        return all_tables
    
    def extract_tabula_tables(self, pdf_path):
        """Extract tables using tabula with optimizations"""
        tables = []
        try:
            # Use faster settings
            tabula_tables = tabula.read_pdf(
                pdf_path, 
                pages='all', 
                multiple_tables=True, 
                silent=True,
                pandas_options={'dtype': str}  # Prevent type inference overhead
            )
            
            for i, table in enumerate(tabula_tables):
                if not table.empty and table.shape[0] > 1 and table.shape[1] > 1:
                    # Quick validation
                    if self.validate_extracted_table_fast(table):
                        tables.append({
                            'method': 'tabula',
                            'table_index': i,
                            'data': table.to_dict('records'),
                            'shape': table.shape,
                            'confidence': 'high'
                        })
        except Exception as e:
            pass
        
        return tables
    
    def extract_camelot_tables(self, pdf_path):
        """Extract tables using camelot with optimizations"""
        tables = []
        try:
            # Use faster lattice method only
            camelot_tables = camelot.read_pdf(pdf_path, pages='all', flavor='lattice')
            
            for i, table in enumerate(camelot_tables):
                if not table.df.empty and table.df.shape[0] > 1 and table.df.shape[1] > 1:
                    accuracy = getattr(table, 'accuracy', 0)
                    if accuracy > 50:
                        tables.append({
                            'method': 'camelot_lattice',
                            'table_index': i,
                            'data': table.df.to_dict('records'),
                            'shape': table.df.shape,
                            'accuracy': accuracy,
                            'confidence': 'high' if accuracy > 80 else 'medium'
                        })
        except Exception as e:
            pass
        
        return tables
    
    def extract_pdfplumber_tables(self, pdf_path):
        """Extract tables using pdfplumber with optimizations"""
        tables = []
        try:
            with pdfplumber.open(pdf_path) as pdf:
                # Process only first 10 pages for speed (adjust as needed)
                pages_to_process = min(10, len(pdf.pages))
                
                for page_num in range(pages_to_process):
                    page = pdf.pages[page_num]
                    page_tables = page.extract_tables()
                    
                    for i, table in enumerate(page_tables):
                        if table and len(table) > 1 and len(table[0]) > 1:
                            if self.validate_pdfplumber_table_fast(table):
                                tables.append({
                                    'method': 'pdfplumber',
                                    'page': page_num + 1,
                                    'table_index': i,
                                    'data': table,
                                    'shape': (len(table), len(table[0]) if table else 0),
                                    'confidence': 'medium'
                                })
        except Exception as e:
            pass
        
        return tables
    
    def validate_extracted_table_fast(self, df):
        """Fast validation of pandas DataFrame"""
        if df.empty or df.shape[0] < 2 or df.shape[1] < 2:
            return False
        
        # Quick check: count non-null values
        non_null_ratio = df.notna().sum().sum() / (df.shape[0] * df.shape[1])
        return non_null_ratio > 0.3
    
    def validate_pdfplumber_table_fast(self, table):
        """Fast validation of pdfplumber table"""
        if not table or len(table) < 2:
            return False
        
        # Quick check: count non-empty cells
        non_empty = sum(1 for row in table for cell in row if cell and cell.strip())
        total = len(table) * len(table[0])
        
        return non_empty / total > 0.3
    
    def extract_docx(self, docx_path):
        """Extract data from DOCX files with optimizations"""
        if not HAS_DOCX:
            print("❌ python-docx not available. Skipping DOCX processing.")
            return None
            
        try:
            print(f"📄 Processing DOCX: {os.path.basename(docx_path)}")
            doc = Document(docx_path)
            
            docx_data = {
                'filename': os.path.basename(docx_path),
                'file_type': 'DOCX',
                'paragraphs': [],
                'tables': [],
                'total_words': 0,
                'total_characters': 0,
                'file_size_mb': round(os.path.getsize(docx_path) / (1024*1024), 2)
            }
            
            # Process paragraphs and tables in parallel
            with ThreadPoolExecutor(max_workers=2) as executor:
                para_future = executor.submit(self.extract_docx_paragraphs, doc)
                table_future = executor.submit(self.extract_docx_tables, doc)
                
                paragraphs, total_words, total_chars = para_future.result()
                tables = table_future.result()
            
            docx_data['paragraphs'] = paragraphs
            docx_data['tables'] = tables
            docx_data['total_words'] = total_words
            docx_data['total_characters'] = total_chars
            
            print(f"  ✅ DOCX processed: {len(docx_data['paragraphs'])} paragraphs, {docx_data['total_words']} words")
            return docx_data
            
        except Exception as e:
            print(f"❌ Error processing DOCX {docx_path}: {str(e)}")
            return None
    
    def extract_docx_paragraphs(self, doc):
        """Extract paragraphs from DOCX"""
        paragraphs = []
        total_words = 0
        total_chars = 0
        
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append({
                    'text': para.text,
                    'style': para.style.name if para.style else 'Normal'
                })
                total_words += len(para.text.split())
                total_chars += len(para.text)
        
        return paragraphs, total_words, total_chars
    
    def extract_docx_tables(self, doc):
        """Extract tables from DOCX"""
        tables = []
        
        for table_idx, table in enumerate(doc.tables):
            table_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(row_data)
            
            if table_data and len(table_data) > 1:
                tables.append({
                    'table_index': table_idx,
                    'data': table_data,
                    'shape': (len(table_data), len(table_data[0]) if table_data else 0)
                })
        
        return tables
    
    def extract_pptx(self, pptx_path):
        """Extract data from PPTX files with optimizations"""
        if not HAS_PPTX:
            print("❌ python-pptx not available. Skipping PPTX processing.")
            return None
            
        try:
            print(f"📄 Processing PPTX: {os.path.basename(pptx_path)}")
            prs = Presentation(pptx_path)
            
            pptx_data = {
                'filename': os.path.basename(pptx_path),
                'file_type': 'PPTX',
                'slides': [],
                'total_slides': len(prs.slides),
                'total_words': 0,
                'total_characters': 0,
                'file_size_mb': round(os.path.getsize(pptx_path) / (1024*1024), 2)
            }
            
            # Process slides in parallel
            with ThreadPoolExecutor(max_workers=self.max_workers) as executor:
                slide_futures = []
                
                for slide_idx, slide in enumerate(prs.slides):
                    future = executor.submit(self.extract_slide_data, slide, slide_idx)
                    slide_futures.append(future)
                
                # Collect results
                for future in slide_futures:
                    slide_data = future.result()
                    pptx_data['slides'].append(slide_data)
                    pptx_data['total_words'] += slide_data['word_count']
                    pptx_data['total_characters'] += slide_data['char_count']
            
            print(f"  ✅ PPTX processed: {pptx_data['total_slides']} slides, {pptx_data['total_words']} words")
            return pptx_data
            
        except Exception as e:
            print(f"❌ Error processing PPTX {pptx_path}: {str(e)}")
            return None
    
    def extract_slide_data(self, slide, slide_idx):
        """Extract data from a single slide"""
        slide_data = {
            'slide_number': slide_idx + 1,
            'text_content': [],
            'tables': [],
            'word_count': 0,
            'char_count': 0
        }
        
        # Extract text and tables from shapes
        for shape in slide.shapes:
            try:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    slide_data['text_content'].append(text)
                    slide_data['word_count'] += len(text.split())
                    slide_data['char_count'] += len(text)
                
                # Extract tables
                if hasattr(shape, "table") and shape.table:
                    table_data = []
                    for row in shape.table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_data.append(row_data)
                    if table_data and len(table_data) > 1:
                        slide_data['tables'].append(table_data)
            except Exception as e:
                continue
        
        return slide_data
    
    def extract_excel(self, excel_path):
        """Extract data from Excel files with optimizations"""
        try:
            print(f"📄 Processing Excel: {os.path.basename(excel_path)}")
            
            excel_data = {
                'filename': os.path.basename(excel_path),
                'file_type': 'Excel',
                'sheets': [],
                'file_size_mb': round(os.path.getsize(excel_path) / (1024*1024), 2)
            }
            
            # Read all sheets with optimizations
            xl_file = pd.ExcelFile(excel_path)
            
            # Process sheets in parallel
            with ThreadPoolExecutor(max_workers=self.max_workers) as executor:
                sheet_futures = []
                
                for sheet_name in xl_file.sheet_names:
                    future = executor.submit(self.extract_excel_sheet, excel_path, sheet_name)
                    sheet_futures.append(future)
                
                # Collect results
                for future in sheet_futures:
                    sheet_data = future.result()
                    if sheet_data:
                        excel_data['sheets'].append(sheet_data)
            
            print(f"  ✅ Excel processed: {len(excel_data['sheets'])} sheets")
            return excel_data
            
        except Exception as e:
            print(f"❌ Error processing Excel {excel_path}: {str(e)}")
            return None
    
    def extract_excel_sheet(self, excel_path, sheet_name):
        """Extract data from a single Excel sheet"""
        try:
            # Read with optimizations
            df = pd.read_excel(
                excel_path, 
                sheet_name=sheet_name,
                nrows=1000,  # Limit rows for speed
                dtype=str    # Prevent type inference
            )
            
            sheet_data = {
                'sheet_name': sheet_name,
                'shape': df.shape,
                'columns': df.columns.tolist(),
                'data': df.head(50).to_dict('records') if not df.empty else [],  # Reduced to 50 rows
                'data_types': df.dtypes.astype(str).to_dict()
            }
            
            return sheet_data
            
        except Exception as e:
            print(f"    ⚠️  Error reading sheet '{sheet_name}': {str(e)}")
            return None
    
    def process_files(self, file_paths):
        """Process all uploaded files with parallel processing"""
        results = {}
        
        print(f"\n🔄 Processing {len(file_paths)} file(s) with {self.max_workers} workers...")
        print("=" * 60)
        
        # Process files in parallel
        with ProcessPoolExecutor(max_workers=self.max_workers) as executor:
            # Submit all files for processing
            future_to_file = {}
            
            for file_path in file_paths:
                file_ext = os.path.splitext(file_path)[1].lower()
                
                if file_ext in self.supported_formats:
                    future = executor.submit(self.process_single_file, file_path)
                    future_to_file[future] = file_path
                else:
                    print(f"❌ Unsupported file format: {file_ext}")
            
            # Collect results as they complete
            for future in future_to_file:
                file_path = future_to_file[future]
                try:
                    result = future.result(timeout=300)  # 5 minute timeout per file
                    if result:
                        results[file_path] = result
                        print(f"✅ Successfully processed: {file_path}")
                    else:
                        print(f"❌ Failed to process: {file_path}")
                except Exception as e:
                    print(f"❌ Error processing {file_path}: {str(e)}")
        
        return results
    
    def process_single_file(self, file_path):
        """Process a single file - used for parallel processing"""
        file_ext = os.path.splitext(file_path)[1].lower()
        
        try:
            if file_ext == '.pdf':
                return self.extract_pdf_comprehensive(file_path)
            elif file_ext == '.docx':
                return self.extract_docx(file_path)
            elif file_ext == '.pptx':
                return self.extract_pptx(file_path)
            elif file_ext in ['.xlsx', '.xls']:
                return self.extract_excel(file_path)
            else:
                return None
                
        except Exception as e:
            print(f"❌ Error processing {file_path}: {str(e)}")
            return None
    
    def save_results(self, results, output_format='json'):
        """Save extraction results with optimizations"""
        if not results:
            print("❌ No results to save.")
            return
        
        try:
            # Save as JSON with optimizations
            with open('extraction_results.json', 'w', encoding='utf-8') as f:
                json.dump(results, f, indent=1, ensure_ascii=False, default=str)  # Reduced indent
            print("✅ Results saved as: extraction_results.json")
            
            # Create summary report and save tables in parallel
            with ThreadPoolExecutor(max_workers=2) as executor:
                summary_future = executor.submit(self.create_summary_report, results)
                csv_future = executor.submit(self.save_tables_as_csv, results)
                
                summary_future.result()
                csv_future.result()
            
            # Check if there are extracted images
            if os.path.exists('extracted_images') and os.listdir('extracted_images'):
                print("🖼️  Images extracted to 'extracted_images' folder")
                
        except Exception as e:
            print(f"❌ Error saving results: {str(e)}")
    
    def save_tables_as_csv(self, results):
        """Save extracted tables as CSV files"""
        table_count = 0
        
        for file_path, data in results.items():
            filename_base = os.path.splitext(data['filename'])[0]
            
            # Save tables from different extraction methods
            if 'extracted_tables' in data:
                for i, table in enumerate(data['extracted_tables']):
                    try:
                        if 'data' in table and table['data']:
                            csv_filename = f"extracted_tables/{filename_base}_table_{i+1}_{table['method']}.csv"
                            
                            if table['method'] in ['tabula', 'camelot_lattice']:
                                # These are already in records format
                                df = pd.DataFrame(table['data'])
                            else:
                                # Convert list of lists to DataFrame
                                df = pd.DataFrame(table['data'])
                            
                            df.to_csv(csv_filename, index=False)
                            table_count += 1
                            
                    except Exception as e:
                        print(f"    ⚠️  Error saving table {i+1}: {str(e)}")
                        continue
        
        if table_count > 0:
            print(f"📊 Saved {table_count} tables as CSV files in 'extracted_tables' folder")
    
    def create_summary_report(self, results):
        """Create a comprehensive summary report"""
        try:
            with open('extraction_summary.txt', 'w', encoding='utf-8') as f:
                f.write("📋 DOCUMENT EXTRACTION SUMMARY REPORT\n")
                f.write("=" * 60 + "\n")
                f.write(f"📅 Generated: {pd.Timestamp.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
                f.write(f"📁 Total Files Processed: {len(results)}\n\n")
                
                # Overall statistics
                total_words = sum(data.get('total_words', 0) for data in results.values())
                total_chars = sum(data.get('total_characters', 0) for data in results.values())
                total_images = sum(data.get('total_images', 0) for data in results.values())
                total_tables = sum(len(data.get('extracted_tables', [])) for data in results.values())
                total_size_mb = sum(data.get('file_size_mb', 0) for data in results.values())
                
                f.write("📊 OVERALL STATISTICS:\n")
                f.write("-" * 30 + "\n")
                f.write(f"📝 Total Words: {total_words:,}\n")
                f.write(f"🔤 Total Characters: {total_chars:,}\n")
                f.write(f"🖼️  Total Images: {total_images:,}\n")
                f.write(f"📊 Total Tables: {total_tables:,}\n")
                f.write(f"💾 Total File Size: {total_size_mb:.2f} MB\n\n")
                
                # File type distribution
                file_types = Counter(data['file_type'] for data in results.values())
                f.write("📈 FILE TYPE DISTRIBUTION:\n")
                f.write("-" * 30 + "\n")
                for file_type, count in file_types.items():
                    f.write(f"{file_type}: {count} files\n")
                f.write("\n")
                
                # Detailed file information
                f.write("📄 DETAILED FILE INFORMATION:\n")
                f.write("=" * 60 + "\n")
                
                for file_path, data in results.items():
                    f.write(f"\n📁 File: {data['filename']}\n")
                    f.write(f"   Type: {data['file_type']}\n")
                    f.write(f"   Size: {data.get('file_size_mb', 0):.2f} MB\n")
                    
                    if data['file_type'] == 'PDF':
                        f.write(f"   Pages: {data.get('page_count', 0)}\n")
                        f.write(f"   Words: {data.get('total_words', 0):,}\n")
                        f.write(f"   Characters: {data.get('total_characters', 0):,}\n")
                        f.write(f"   Images: {data.get('total_images', 0)}\n")
                        f.write(f"   Tables: {len(data.get('extracted_tables', []))}\n")
                        
                        # Font information
                        fonts = data.get('fonts_used', [])
                        if fonts:
                            f.write(f"   Fonts Used: {', '.join(fonts[:5])}")
                            if len(fonts) > 5:
                                f.write(f" ... (+{len(fonts)-5} more)")
                            f.write("\n")
                        
                        # Metadata
                        metadata = data.get('metadata', {})
                        if metadata:
                            f.write("   Metadata:\n")
                            for key, value in metadata.items():
                                f.write(f"     {key}: {value}\n")
                    
                    elif data['file_type'] == 'DOCX':
                        f.write(f"   Paragraphs: {len(data.get('paragraphs', []))}\n")
                        f.write(f"   Words: {data.get('total_words', 0):,}\n")
                        f.write(f"   Characters: {data.get('total_characters', 0):,}\n")
                        f.write(f"   Tables: {len(data.get('tables', []))}\n")
                    
                    elif data['file_type'] == 'PPTX':
                        f.write(f"   Slides: {data.get('total_slides', 0)}\n")
                        f.write(f"   Words: {data.get('total_words', 0):,}\n")
                        f.write(f"   Characters: {data.get('total_characters', 0):,}\n")
                        
                        # Count tables across all slides
                        total_slide_tables = sum(len(slide.get('tables', [])) for slide in data.get('slides', []))
                        f.write(f"   Tables: {total_slide_tables}\n")
                    
                    elif data['file_type'] == 'Excel':
                        f.write(f"   Sheets: {len(data.get('sheets', []))}\n")
                        
                        # Sheet details
                        for sheet in data.get('sheets', []):
                            f.write(f"     - {sheet['sheet_name']}: {sheet['shape'][0]} rows x {sheet['shape'][1]} cols\n")
                
                # Table extraction summary
                if total_tables > 0:
                    f.write(f"\n📊 TABLE EXTRACTION SUMMARY:\n")
                    f.write("-" * 40 + "\n")
                    
                    method_counts = Counter()
                    confidence_counts = Counter()
                    
                    for data in results.values():
                        for table in data.get('extracted_tables', []):
                            method_counts[table.get('method', 'unknown')] += 1
                            confidence_counts[table.get('confidence', 'unknown')] += 1
                    
                    f.write("Extraction Methods:\n")
                    for method, count in method_counts.items():
                        f.write(f"  {method}: {count} tables\n")
                    
                    f.write("\nConfidence Levels:\n")
                    for confidence, count in confidence_counts.items():
                        f.write(f"  {confidence}: {count} tables\n")
                
                # Image extraction summary
                if total_images > 0:
                    f.write(f"\n🖼️  IMAGE EXTRACTION SUMMARY:\n")
                    f.write("-" * 40 + "\n")
                    
                    for data in results.values():
                        if data.get('total_images', 0) > 0:
                            f.write(f"From {data['filename']}: {data['total_images']} images\n")
                            
                            # Image details from pages
                            for page in data.get('pages', []):
                                if page.get('images'):
                                    f.write(f"  Page {page['page_number']}: {len(page['images'])} images\n")
                
                # Processing warnings and recommendations
                f.write(f"\n⚠️  PROCESSING NOTES:\n")
                f.write("-" * 30 + "\n")
                
                # Check for potential issues
                large_files = [data for data in results.values() if data.get('file_size_mb', 0) > 50]
                if large_files:
                    f.write(f"📁 Large files detected ({len(large_files)} files > 50MB)\n")
                
                # Check for files with many images
                image_heavy_files = [data for data in results.values() if data.get('total_images', 0) > 20]
                if image_heavy_files:
                    f.write(f"🖼️  Image-heavy files detected ({len(image_heavy_files)} files > 20 images)\n")
                
                # Check for missing libraries
                missing_libs = []
                if not HAS_PDFPLUMBER:
                    missing_libs.append("pdfplumber")
                if not HAS_TABULA:
                    missing_libs.append("tabula-py")
                if not HAS_CAMELOT:
                    missing_libs.append("camelot-py")
                if not HAS_CV2:
                    missing_libs.append("opencv-python")
                
                if missing_libs:
                    f.write(f"📦 Missing optional libraries: {', '.join(missing_libs)}\n")
                    f.write("   Install these for enhanced extraction capabilities\n")
                
                f.write(f"\n✅ Report generated successfully!\n")
                f.write("   Files created:\n")
                f.write("   - extraction_results.json (full data)\n")
                f.write("   - extraction_summary.txt (this report)\n")
                f.write("   - extracted_images/ (folder with images)\n")
                f.write("   - extracted_tables/ (folder with CSV tables)\n")
            
            print("✅ Summary report saved as: extraction_summary.txt")
            
        except Exception as e:
            print(f"❌ Error creating summary report: {str(e)}")
    
    def generate_extraction_stats(self, results):
        """Generate detailed extraction statistics"""
        stats = {
            'total_files': len(results),
            'file_types': Counter(data['file_type'] for data in results.values()),
            'total_words': sum(data.get('total_words', 0) for data in results.values()),
            'total_characters': sum(data.get('total_characters', 0) for data in results.values()),
            'total_images': sum(data.get('total_images', 0) for data in results.values()),
            'total_tables': sum(len(data.get('extracted_tables', [])) for data in results.values()),
            'total_size_mb': sum(data.get('file_size_mb', 0) for data in results.values()),
            'processing_time': getattr(self, 'processing_time', 0)
        }
        
        # Add file-specific stats
        stats['pdf_stats'] = {
            'count': sum(1 for data in results.values() if data['file_type'] == 'PDF'),
            'total_pages': sum(data.get('page_count', 0) for data in results.values() if data['file_type'] == 'PDF'),
            'fonts_used': list(set(font for data in results.values() if data['file_type'] == 'PDF' for font in data.get('fonts_used', [])))
        }
        
        return stats
    
    def cleanup_temp_files(self):
        """Clean up temporary files and directories"""
        try:
            # Remove empty directories
            for folder in ['extracted_images', 'extracted_tables']:
                if os.path.exists(folder) and not os.listdir(folder):
                    os.rmdir(folder)
                    print(f"🧹 Removed empty folder: {folder}")
        except Exception as e:
            print(f"⚠️  Error during cleanup: {str(e)}")

def main():
    """Main function to run the document extractor"""
    print("🚀 Comprehensive Document Extractor")
    print("=" * 60)
    print("Supported formats: PDF, DOCX, PPTX, XLSX, XLS")
    print()
    
    # Initialize the extractor
    extractor = ComprehensiveDocumentExtractor()
    
    # Get file paths from user
    print("📁 Please provide file paths to process:")
    print("   - Enter file paths separated by commas")
    print("   - Or drag and drop files into the terminal")
    print("   - Press Enter when done")
    print()
    
    file_paths = []
    while True:
        user_input = input("File path (or 'done' to finish): ").strip()
        
        if user_input.lower() in ['done', 'finish', 'exit', '']:
            break
        
        # Handle multiple files separated by commas
        if ',' in user_input:
            paths = [path.strip().strip('"\'') for path in user_input.split(',')]
            file_paths.extend(paths)
        else:
            file_paths.append(user_input.strip('"\''))
    
    if not file_paths:
        print("❌ No files provided. Exiting.")
        return
    
    # Validate file paths
    valid_paths = []
    for path in file_paths:
        if os.path.exists(path):
            ext = os.path.splitext(path)[1].lower()
            if ext in extractor.supported_formats:
                valid_paths.append(path)
            else:
                print(f"⚠️  Skipping unsupported file: {path}")
        else:
            print(f"⚠️  File not found: {path}")
    
    if not valid_paths:
        print("❌ No valid files to process. Exiting.")
        return
    
    # Record start time
    start_time = pd.Timestamp.now()
    
    # Process files
    try:
        results = extractor.process_files(valid_paths)
        
        if results:
            # Record processing time
            extractor.processing_time = (pd.Timestamp.now() - start_time).total_seconds()
            
            # Save results
            extractor.save_results(results)
            
            # Generate and display final statistics
            stats = extractor.generate_extraction_stats(results)
            
            print("\n" + "=" * 60)
            print("📊 FINAL EXTRACTION STATISTICS")
            print("=" * 60)
            print(f"✅ Successfully processed: {stats['total_files']} files")
            print(f"📝 Total words extracted: {stats['total_words']:,}")
            print(f"🔤 Total characters: {stats['total_characters']:,}")
            print(f"🖼️  Total images: {stats['total_images']:,}")
            print(f"📊 Total tables: {stats['total_tables']:,}")
            print(f"💾 Total file size: {stats['total_size_mb']:.2f} MB")
            print(f"⏱️  Processing time: {stats['processing_time']:.2f} seconds")
            
            if stats['pdf_stats']['count'] > 0:
                print(f"📄 PDF files: {stats['pdf_stats']['count']} ({stats['pdf_stats']['total_pages']} pages)")
            
            print("\n📁 Output files created:")
            print("   - extraction_results.json")
            print("   - extraction_summary.txt")
            if stats['total_images'] > 0:
                print("   - extracted_images/ (folder)")
            if stats['total_tables'] > 0:
                print("   - extracted_tables/ (folder)")
            
            # Cleanup
            extractor.cleanup_temp_files()
            
            print("\n🎉 Document extraction completed successfully!")
            
        else:
            print("❌ No files were successfully processed.")
            
    except KeyboardInterrupt:
        print("\n⏹️  Process interrupted by user.")
        extractor.cleanup_temp_files()
    except Exception as e:
        print(f"\n❌ Unexpected error: {str(e)}")
        extractor.cleanup_temp_files()

if __name__ == "__main__":
    main()
