import streamlit as st
import fitz  # PyMuPDF
import io
import os
import json
import pandas as pd
from PIL import Image
import pytesseract
from collections import Counter
import zipfile
import tempfile
import pdfplumber
import tabula
import numpy as np
from docx import Document
from pptx import Presentation
import openpyxl
import warnings
import base64
from datetime import datetime
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots

warnings.filterwarnings('ignore')

# Page configuration
st.set_page_config(
    page_title="Document Intelligence Extractor",
    page_icon="📄",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS for technical styling
st.markdown("""
<style>
    .main-header {
        font-size: 3rem;
        font-weight: 700;
        background: linear-gradient(90deg, #667eea 0%, #764ba2 100%);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        text-align: center;
        margin-bottom: 1rem;
    }
    
    .tech-card {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        padding: 1rem;
        border-radius: 10px;
        color: white;
        margin: 0.5rem 0;
    }
    
    .metric-card {
        background: linear-gradient(135deg, #f093fb 0%, #f5576c 100%);
        padding: 1rem;
        border-radius: 10px;
        color: white;
        text-align: center;
        margin: 0.5rem;
    }
    
    .success-box {
        background: linear-gradient(135deg, #4facfe 0%, #00f2fe 100%);
        padding: 1rem;
        border-radius: 10px;
        color: white;
        margin: 1rem 0;
    }
    
    .warning-box {
        background: linear-gradient(135deg, #fa709a 0%, #fee140 100%);
        padding: 1rem;
        border-radius: 10px;
        color: white;
        margin: 1rem 0;
    }
    
    .stProgress > div > div > div > div {
        background: linear-gradient(90deg, #667eea 0%, #764ba2 100%);
    }
    
    .sidebar-content {
        background: linear-gradient(180deg, #667eea 0%, #764ba2 100%);
        padding: 1rem;
        border-radius: 10px;
        color: white;
        margin: 1rem 0;
    }
</style>
""", unsafe_allow_html=True)

class StreamlitDocumentExtractor:
    def __init__(self):
        self.supported_formats = ['.pdf', '.docx', '.pptx', '.xlsx', '.xls']
        self.setup_session_state()
    
    def setup_session_state(self):
        """Initialize session state variables"""
        if 'extraction_results' not in st.session_state:
            st.session_state.extraction_results = {}
        if 'processing_complete' not in st.session_state:
            st.session_state.processing_complete = False
        if 'uploaded_files' not in st.session_state:
            st.session_state.uploaded_files = []
    
    def create_download_link(self, data, filename, link_text):
        """Create a download link for data"""
        if isinstance(data, dict):
            data = json.dumps(data, indent=2, ensure_ascii=False, default=str)
        
        b64 = base64.b64encode(data.encode()).decode()
        href = f'<a href="data:file/json;base64,{b64}" download="{filename}" style="text-decoration: none; background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); color: white; padding: 0.5rem 1rem; border-radius: 5px; margin: 0.25rem;">{link_text}</a>'
        return href
    
    def extract_pdf_comprehensive(self, uploaded_file):
        """Extract comprehensive data from PDF with progress tracking"""
        try:
            # Save uploaded file temporarily
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_file:
                tmp_file.write(uploaded_file.read())
                tmp_path = tmp_file.name
            
            doc = fitz.open(tmp_path)
            
            pdf_data = {
                'filename': uploaded_file.name,
                'file_type': 'PDF',
                'metadata': {},
                'pages': [],
                'total_images': 0,
                'total_tables': 0,
                'fonts_used': [],
                'total_words': 0,
                'total_characters': 0,
                'page_count': len(doc),
                'file_size_mb': round(len(uploaded_file.getvalue()) / (1024*1024), 2)
            }
            
            # Extract metadata safely
            try:
                metadata = doc.metadata
                pdf_data['metadata'] = {k: v for k, v in metadata.items() if v}
            except:
                pdf_data['metadata'] = {}
            
            # Create progress bar
            progress_bar = st.progress(0)
            status_text = st.empty()
            
            # Process each page
            for page_num in range(len(doc)):
                progress = (page_num + 1) / len(doc)
                progress_bar.progress(progress)
                status_text.text(f"Processing page {page_num + 1}/{len(doc)}")
                
                try:
                    page = doc[page_num]
                    page_data = self.extract_page_data(page, page_num + 1)
                    pdf_data['pages'].append(page_data)
                    
                    # Accumulate totals
                    pdf_data['total_images'] += len(page_data['images'])
                    pdf_data['total_tables'] += len(page_data['tables'])
                    pdf_data['total_words'] += page_data['word_count']
                    pdf_data['total_characters'] += page_data['char_count']
                    
                except Exception as e:
                    st.warning(f"Error on page {page_num + 1}: {str(e)}")
                    continue
            
            # Extract unique fonts
            all_fonts = []
            for page in pdf_data['pages']:
                all_fonts.extend(page.get('fonts', []))
            pdf_data['fonts_used'] = list(set(all_fonts))
            
            # Extract tables
            status_text.text("Extracting tables...")
            pdf_data['extracted_tables'] = self.extract_tables_multiple_methods(tmp_path)
            
            doc.close()
            os.unlink(tmp_path)  # Clean up temp file
            
            progress_bar.progress(1.0)
            status_text.text("✅ PDF processing complete!")
            
            return pdf_data
            
        except Exception as e:
            st.error(f"Error processing PDF: {str(e)}")
            return None
    
    def extract_page_data(self, page, page_num):
        """Extract comprehensive data from a single page"""
        page_data = {
            'page_number': page_num,
            'text': '',
            'formatted_text': [],
            'images': [],
            'tables': [],
            'fonts': [],
            'word_count': 0,
            'char_count': 0,
            'page_size': {
                'width': round(page.rect.width, 2),
                'height': round(page.rect.height, 2)
            },
            'rotation': page.rotation
        }
        
        try:
            # Extract text with formatting
            text_dict = page.get_text("dict")
            page_text = ""
            fonts_on_page = []
            
            for block in text_dict.get("blocks", []):
                if "lines" in block:
                    for line in block["lines"]:
                        for span in line["spans"]:
                            text = span.get("text", "")
                            page_text += text
                            
                            page_data['formatted_text'].append({
                                'text': text,
                                'font': span.get('font', 'Unknown'),
                                'size': round(span.get('size', 0), 2),
                                'flags': span.get('flags', 0),
                                'color': span.get('color', 0),
                                'bbox': span.get('bbox', [0, 0, 0, 0])
                            })
                            
                            fonts_on_page.append(span.get('font', 'Unknown'))
            
            page_data['text'] = page_text
            page_data['fonts'] = list(set(fonts_on_page))
            page_data['word_count'] = len(page_text.split())
            page_data['char_count'] = len(page_text)
            
            # Extract images
            page_data['images'] = self.extract_images_from_page(page, page_num)
            
            # Extract tables
            page_data['tables'] = self.extract_tables_from_page(page)
            
        except Exception as e:
            st.warning(f"Error extracting from page {page_num}: {str(e)}")
        
        return page_data
    
    def extract_images_from_page(self, page, page_num):
        """Extract images from a page"""
        images = []
        
        try:
            image_list = page.get_images()
            
            for img_index, img in enumerate(image_list):
                try:
                    xref = img[0]
                    pix = fitz.Pixmap(page.parent, xref)
                    
                    if pix.n - pix.alpha < 4:  # GRAY or RGB
                        img_data = pix.tobytes("png")
                        
                        img_info = {
                            'page': page_num,
                            'index': img_index + 1,
                            'width': pix.width,
                            'height': pix.height,
                            'colorspace': pix.colorspace.name if pix.colorspace else 'Unknown',
                            'size_bytes': len(img_data),
                            'data': base64.b64encode(img_data).decode()
                        }
                        
                        images.append(img_info)
                    
                    pix = None
                    
                except Exception as e:
                    continue
                    
        except Exception as e:
            pass
        
        return images
    
    def extract_tables_from_page(self, page):
        """Basic table detection from page"""
        tables = []
        try:
            text_dict = page.get_text("dict")
            
            for block in text_dict.get("blocks", []):
                if "lines" in block:
                    lines = block["lines"]
                    if len(lines) > 2:
                        table_data = []
                        for line in lines:
                            row_text = ""
                            for span in line["spans"]:
                                row_text += span.get("text", "") + " "
                            table_data.append(row_text.strip())
                        
                        if len(table_data) > 2:
                            tables.append({
                                'bbox': block.get('bbox', [0, 0, 0, 0]),
                                'data': table_data,
                                'method': 'basic_detection'
                            })
            
        except Exception as e:
            pass
        
        return tables
    
    def extract_tables_multiple_methods(self, pdf_path):
        """Extract tables using multiple libraries"""
        all_tables = []
        
        # Method 1: Using tabula-py
        try:
            tabula_tables = tabula.read_pdf(pdf_path, pages='all', multiple_tables=True, silent=True)
            for i, table in enumerate(tabula_tables):
                if not table.empty:
                    all_tables.append({
                        'method': 'tabula',
                        'table_index': i,
                        'data': table.to_dict('records'),
                        'shape': table.shape,
                        'columns': table.columns.tolist()
                    })
        except Exception as e:
            pass
        
        # Method 2: Using pdfplumber
        try:
            table_count = 0
            with pdfplumber.open(pdf_path) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    tables = page.extract_tables()
                    for i, table in enumerate(tables):
                        if table and len(table) > 1:
                            all_tables.append({
                                'method': 'pdfplumber',
                                'page': page_num + 1,
                                'table_index': i,
                                'data': table,
                                'shape': (len(table), len(table[0]) if table else 0)
                            })
                            table_count += 1
        except Exception as e:
            pass
        
        return all_tables
    
    def extract_docx(self, uploaded_file):
        """Extract data from DOCX files"""
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as tmp_file:
                tmp_file.write(uploaded_file.read())
                tmp_path = tmp_file.name
            
            doc = Document(tmp_path)
            
            docx_data = {
                'filename': uploaded_file.name,
                'file_type': 'DOCX',
                'paragraphs': [],
                'tables': [],
                'total_words': 0,
                'total_characters': 0,
                'file_size_mb': round(len(uploaded_file.getvalue()) / (1024*1024), 2)
            }
            
            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    docx_data['paragraphs'].append({
                        'text': para.text,
                        'style': para.style.name if para.style else 'Normal'
                    })
                    docx_data['total_words'] += len(para.text.split())
                    docx_data['total_characters'] += len(para.text)
            
            # Extract tables
            for table_idx, table in enumerate(doc.tables):
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                
                if table_data:
                    docx_data['tables'].append({
                        'table_index': table_idx,
                        'data': table_data,
                        'shape': (len(table_data), len(table_data[0]) if table_data else 0)
                    })
            
            os.unlink(tmp_path)
            return docx_data
            
        except Exception as e:
            st.error(f"Error processing DOCX: {str(e)}")
            return None
    
    def extract_pptx(self, uploaded_file):
        """Extract data from PPTX files"""
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
                tmp_file.write(uploaded_file.read())
                tmp_path = tmp_file.name
            
            prs = Presentation(tmp_path)
            
            pptx_data = {
                'filename': uploaded_file.name,
                'file_type': 'PPTX',
                'slides': [],
                'total_slides': len(prs.slides),
                'total_words': 0,
                'total_characters': 0,
                'file_size_mb': round(len(uploaded_file.getvalue()) / (1024*1024), 2)
            }
            
            for slide_idx, slide in enumerate(prs.slides):
                slide_data = {
                    'slide_number': slide_idx + 1,
                    'text_content': [],
                    'tables': []
                }
                
                # Extract text from shapes
                for shape in slide.shapes:
                    try:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_data['text_content'].append(shape.text.strip())
                            pptx_data['total_words'] += len(shape.text.split())
                            pptx_data['total_characters'] += len(shape.text)
                        
                        # Extract tables
                        if hasattr(shape, "table") and shape.table:
                            table_data = []
                            for row in shape.table.rows:
                                row_data = [cell.text.strip() for cell in row.cells]
                                table_data.append(row_data)
                            if table_data:
                                slide_data['tables'].append(table_data)
                    except Exception as e:
                        continue
                
                pptx_data['slides'].append(slide_data)
            
            os.unlink(tmp_path)
            return pptx_data
            
        except Exception as e:
            st.error(f"Error processing PPTX: {str(e)}")
            return None
    
    def extract_excel(self, uploaded_file):
        """Extract data from Excel files"""
        try:
            excel_data = {
                'filename': uploaded_file.name,
                'file_type': 'Excel',
                'sheets': [],
                'file_size_mb': round(len(uploaded_file.getvalue()) / (1024*1024), 2)
            }
            
            # Read all sheets
            xl_file = pd.ExcelFile(uploaded_file)
            
            for sheet_name in xl_file.sheet_names:
                try:
                    df = pd.read_excel(uploaded_file, sheet_name=sheet_name)
                    
                    sheet_data = {
                        'sheet_name': sheet_name,
                        'shape': df.shape,
                        'columns': df.columns.tolist(),
                        'data': df.head(100).to_dict('records') if not df.empty else [],
                        'summary': df.describe(include='all').to_dict() if not df.empty else {},
                        'data_types': df.dtypes.astype(str).to_dict()
                    }
                    
                    excel_data['sheets'].append(sheet_data)
                    
                except Exception as e:
                    st.warning(f"Error reading sheet '{sheet_name}': {str(e)}")
                    continue
            
            return excel_data
            
        except Exception as e:
            st.error(f"Error processing Excel: {str(e)}")
            return None
    
    def generate_summary_report(self, results):
        """Generate a comprehensive summary report"""
        report = {
            'extraction_timestamp': datetime.now().isoformat(),
            'total_files_processed': len(results),
            'file_types': {},
            'summary_statistics': {
                'total_words': 0,
                'total_characters': 0,
                'total_images': 0,
                'total_tables': 0,
                'total_size_mb': 0
            },
            'file_details': []
        }
        
        for filename, data in results.items():
            # File type distribution
            file_type = data.get('file_type', 'Unknown')
            report['file_types'][file_type] = report['file_types'].get(file_type, 0) + 1
            
            # Accumulate statistics
            report['summary_statistics']['total_words'] += data.get('total_words', 0)
            report['summary_statistics']['total_characters'] += data.get('total_characters', 0)
            report['summary_statistics']['total_images'] += data.get('total_images', 0)
            report['summary_statistics']['total_size_mb'] += data.get('file_size_mb', 0)
            
            # Count tables from different sources
            tables_count = len(data.get('extracted_tables', []))
            tables_count += len(data.get('tables', []))
            if data.get('file_type') == 'Excel':
                for sheet in data.get('sheets', []):
                    if sheet.get('shape', [0, 0])[0] > 0:
                        tables_count += 1
            report['summary_statistics']['total_tables'] += tables_count
            
            # File details
            file_detail = {
                'filename': filename,
                'file_type': file_type,
                'file_size_mb': data.get('file_size_mb', 0),
                'content_summary': {}
            }
            
            if file_type == 'PDF':
                file_detail['content_summary'] = {
                    'pages': data.get('page_count', 0),
                    'words': data.get('total_words', 0),
                    'images': data.get('total_images', 0),
                    'tables': len(data.get('extracted_tables', [])),
                    'fonts': len(data.get('fonts_used', []))
                }
            elif file_type == 'DOCX':
                file_detail['content_summary'] = {
                    'paragraphs': len(data.get('paragraphs', [])),
                    'words': data.get('total_words', 0),
                    'tables': len(data.get('tables', []))
                }
            elif file_type == 'PPTX':
                file_detail['content_summary'] = {
                    'slides': data.get('total_slides', 0),
                    'words': data.get('total_words', 0)
                }
            elif file_type == 'Excel':
                file_detail['content_summary'] = {
                    'sheets': len(data.get('sheets', [])),
                    'total_rows': sum(sheet.get('shape', [0, 0])[0] for sheet in data.get('sheets', [])),
                    'total_columns': sum(sheet.get('shape', [0, 0])[1] for sheet in data.get('sheets', []))
                }
            
            report['file_details'].append(file_detail)
        
        return report

def main():
    # Header
    st.markdown('<h1 class="main-header">📄 Document Intelligence Extractor</h1>', unsafe_allow_html=True)
    
    # Sidebar
    with st.sidebar:
        st.markdown("""
        <div class="sidebar-content">
            <h3>🚀 Features</h3>
            <ul>
                <li>🔍 Text Extraction</li>
                <li>🖼️ Image Extraction</li>
                <li>📊 Table Detection</li>
                <li>📋 Metadata Analysis</li>
                <li>🎨 Font Analysis</li>
                <li>📈 Analytics Dashboard</li>
            </ul>
        </div>
        """, unsafe_allow_html=True)
        
        st.markdown("""
        <div class="sidebar-content">
            <h3>📁 Supported Formats</h3>
            <ul>
                <li>📄 PDF (.pdf)</li>
                <li>📝 Word (.docx)</li>
                <li>📊 Excel (.xlsx, .xls)</li>
                <li>🎯 PowerPoint (.pptx)</li>
            </ul>
        </div>
        """, unsafe_allow_html=True)
        
        # Clear results button
        if st.button("🗑️ Clear Results", use_container_width=True):
            st.session_state.extraction_results = {}
            st.session_state.processing_complete = False
            st.session_state.uploaded_files = []
            st.rerun()
    
    extractor = StreamlitDocumentExtractor()
    
    # File upload section
    st.markdown("### 📤 Upload Documents")
    uploaded_files = st.file_uploader(
        "Choose files to extract data from",
        type=['pdf', 'docx', 'pptx', 'xlsx', 'xls'],
        accept_multiple_files=True,
        help="Upload one or more documents for comprehensive analysis"
    )
    
    if uploaded_files:
        st.markdown(f"<div class='success-box'>✅ {len(uploaded_files)} file(s) uploaded successfully!</div>", unsafe_allow_html=True)
        
        # Display uploaded files info
        col1, col2, col3 = st.columns(3)
        
        with col1:
            st.markdown(f"<div class='metric-card'><h3>{len(uploaded_files)}</h3><p>Files Uploaded</p></div>", unsafe_allow_html=True)
        
        with col2:
            total_size = sum(len(f.getvalue()) for f in uploaded_files) / (1024*1024)
            st.markdown(f"<div class='metric-card'><h3>{total_size:.2f} MB</h3><p>Total Size</p></div>", unsafe_allow_html=True)
        
        with col3:
            formats = list(set(f.name.split('.')[-1].upper() for f in uploaded_files))
            st.markdown(f"<div class='metric-card'><h3>{len(formats)}</h3><p>File Types</p></div>", unsafe_allow_html=True)
        
        # Process files button
        if st.button("🚀 Start Extraction", type="primary", use_container_width=True):
            results = {}
            
            for i, uploaded_file in enumerate(uploaded_files):
                st.markdown(f"### 📊 Processing: {uploaded_file.name}")
                
                file_ext = uploaded_file.name.split('.')[-1].lower()
                
                try:
                    if file_ext == 'pdf':
                        result = extractor.extract_pdf_comprehensive(uploaded_file)
                    elif file_ext == 'docx':
                        result = extractor.extract_docx(uploaded_file)
                    elif file_ext == 'pptx':
                        result = extractor.extract_pptx(uploaded_file)
                    elif file_ext in ['xlsx', 'xls']:
                        result = extractor.extract_excel(uploaded_file)
                    else:
                        st.error(f"Unsupported file format: {file_ext}")
                        continue
                    
                    if result:
                        results[uploaded_file.name] = result
                        st.success(f"✅ Successfully processed: {uploaded_file.name}")
                    else:
                        st.error(f"❌ Failed to process: {uploaded_file.name}")
                        
                except Exception as e:
                    st.error(f"❌ Error processing {uploaded_file.name}: {str(e)}")
                    continue
            
            if results:
                st.session_state.extraction_results = results
                st.session_state.processing_complete = True
                
                # Display results
                display_results(results)
    
    # Display existing results if available
    elif st.session_state.processing_complete and st.session_state.extraction_results:
        display_results(st.session_state.extraction_results)

def display_results(results):
    """Display extraction results with analytics"""
    st.markdown("## 📊 Extraction Results & Analytics")
    
    # Overview metrics
    st.markdown("### 📈 Overview")
    
    col1, col2, col3, col4 = st.columns(4)
    
    total_files = len(results)
    total_words = sum(data.get('total_words', 0) for data in results.values())
    total_images = sum(data.get('total_images', 0) for data in results.values())
    total_tables = sum(len(data.get('extracted_tables', [])) + len(data.get('tables', [])) for data in results.values())
    
    with col1:
        st.metric("📁 Total Files", total_files)
    with col2:
        st.metric("📝 Total Words", f"{total_words:,}")
    with col3:
        st.metric("🖼️ Total Images", total_images)
    with col4:
        st.metric("📊 Total Tables", total_tables)
    
    # File type distribution chart
    file_types = [data.get('file_type', 'Unknown') for data in results.values()]
    file_type_counts = pd.Series(file_types).value_counts()
    
    fig_pie = px.pie(
        values=file_type_counts.values,
        names=file_type_counts.index,
        title="File Type Distribution",
        color_discrete_sequence=px.colors.qualitative.Set3
    )
    st.plotly_chart(fig_pie, use_container_width=True)
    
    # Word count comparison
    if total_words > 0:
        word_data = []
        for filename, data in results.items():
            word_data.append({
                'File': filename,
                'Words': data.get('total_words', 0),
                'Type': data.get('file_type', 'Unknown')
            })
        
        df_words = pd.DataFrame(word_data)
        fig_bar = px.bar(
            df_words,
            x='File',
            y='Words',
            color='Type',
            title="Word Count by File",
            color_discrete_sequence=px.colors.qualitative.Set2
        )
        fig_bar.update_xaxis(tickangle=45)
        st.plotly_chart(fig_bar, use_container_width=True)
    
    # Detailed results for each file
    st.markdown("### 📋 Detailed Results")
    
    for filename, data in results.items():
        with st.expander(f"📄 {filename} - {data.get('file_type', 'Unknown')}"):
            
            # File summary
            col1, col2, col3 = st.columns(3)
            
            with col1:
                st.markdown(f"**📊 File Size:** {data.get('file_size_mb', 0)} MB")
                if data.get('file_type') == 'PDF':
                    st.markdown(f"**📃 Pages:** {data.get('page_count', 0)}")
                elif data.get('file_type') == 'PPTX':
                    st.markdown(f"**🎯 Slides:** {data.get('total_slides', 0)}")
                elif data.get('file_type') == 'Excel':
                    st.markdown(f"**📊 Sheets:** {len(data.get('sheets', []))}")
            
            with col2:
                st.markdown(f"**📝 Words:** {data.get('total_words', 0):,}")
                st.markdown(f"**🔤 Characters:** {data.get('total_characters', 0):,}")
                st.markdown(f"**📋 Tables:** {len(data.get('extracted_tables', [])) + len(data.get('tables', []))}")

            with col3:
                st.markdown(f"**🖼️ Images:** {data.get('total_images', 0)}")
                if data.get('file_type') == 'PDF':
                    st.markdown(f"**🔠 Fonts Used:** {len(data.get('fonts_used', []))}")

            # Metadata (for PDF)
            if data.get('file_type') == 'PDF' and data.get('metadata'):
                st.markdown("#### 🧾 Metadata")
                for key, value in data['metadata'].items():
                    st.markdown(f"- **{key}**: {value}")

            # Download JSON
            st.markdown("#### 📥 Download Data")
            st.markdown(extractor.create_download_link(data, f"{filename}_extraction.json", "💾 Download JSON"), unsafe_allow_html=True)

            # Extracted Tables
            if 'extracted_tables' in data and data['extracted_tables']:
                st.markdown("#### 📊 Extracted Tables")
                for i, table in enumerate(data['extracted_tables']):
                    try:
                        st.markdown(f"**Table {i+1}** (method: {table['method']})")
                        df = pd.DataFrame(table['data'])
                        st.dataframe(df)
                    except Exception as e:
                        st.warning(f"⚠️ Error displaying table {i+1}: {str(e)}")

            # Extracted Images
            if 'images' in data and data['images']:
                st.markdown("#### 🖼️ Extracted Images")
                img_cols = st.columns(4)
                for idx, img in enumerate(data['images']):
                    try:
                        img_data = base64.b64decode(img['data'])
                        img_cols[idx % 4].image(Image.open(io.BytesIO(img_data)), caption=f"Page {img.get('page', '?')}, Img {img.get('index', '?')}", use_column_width=True)
                    except Exception as e:
                        continue