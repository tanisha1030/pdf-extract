import concurrent.futures
import os
import json
import pandas as pd
from PIL import Image
import pytesseract
from collections import Counter
import zipfile
import tempfile
import warnings
import re
import fitz  # PyMuPDF
import gc  # Add garbage collection
import psutil  # Add memory monitoring

# Optional libraries
try:
    import pdfplumber
    HAS_PDFPLUMBER = True
except ImportError:
    HAS_PDFPLUMBER = False

try:
    import tabula
    HAS_TABULA = True
except ImportError:
    HAS_TABULA = False

try:
    import camelot
    HAS_CAMELOT = True
except ImportError:
    HAS_CAMELOT = False

try:
    import cv2
    import numpy as np
    HAS_CV2 = True
except ImportError:
    HAS_CV2 = False

try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import openpyxl
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False


class ComprehensiveDocumentExtractor:
    def __init__(self, max_workers=None, memory_limit_mb=1024, extract_formatting=False):
        self.extracted_data = {}
        self.supported_formats = ['.pdf', '.docx', '.pptx', '.xlsx', '.xls']
        self.memory_limit_mb = memory_limit_mb
        self.extract_formatting = extract_formatting  # Option to disable detailed formatting for large files
        self.max_workers = max_workers or min(8, os.cpu_count() or 4)  # Adaptive worker count
        self.setup_directories()

    def setup_directories(self):
        """Create necessary directories"""
        os.makedirs('extracted_images', exist_ok=True)
        os.makedirs('extracted_tables', exist_ok=True)

    def check_memory_usage(self):
        """Check current memory usage"""
        process = psutil.Process(os.getpid())
        memory_mb = process.memory_info().rss / 1024 / 1024
        return memory_mb

    def is_memory_limit_exceeded(self):
        """Check if memory limit is exceeded"""
        return self.check_memory_usage() > self.memory_limit_mb

    def extract_pdf_comprehensive(self, pdf_path):
        """Extract comprehensive data from PDF with optimized memory usage"""
        try:
            print(f"📄 Processing PDF: {os.path.basename(pdf_path)}")
            doc = fitz.open(pdf_path)
            
            # Check file size and adjust processing accordingly
            file_size_mb = os.path.getsize(pdf_path) / (1024*1024)
            page_count = len(doc)
            print(f"  📊 File size: {file_size_mb:.2f} MB, Pages: {page_count}")
            
            # Adjust processing based on file size
            if file_size_mb > 100:  # Large file optimizations
                print("  🔧 Large file detected - applying optimizations...")
                self.extract_formatting = False  # Disable detailed formatting
                batch_size = max(1, min(10, 50 // int(file_size_mb / 10)))  # Adaptive batch size
            else:
                batch_size = page_count
            
            pdf_data = {
                'filename': os.path.basename(pdf_path),
                'file_type': 'PDF',
                'metadata': {},
                'pages': [],
                'total_images': 0,
                'total_tables': 0,
                'fonts_used': [],
                'total_words': 0,
                'total_characters': 0,
                'page_count': page_count,
                'file_size_mb': round(file_size_mb, 2),
                'processing_mode': 'optimized' if file_size_mb > 100 else 'standard'
            }

            # Extract metadata safely
            try:
                metadata = doc.metadata
                if metadata:
                    pdf_data['metadata'] = {k: v for k, v in metadata.items() if v}
            except:
                pdf_data['metadata'] = {}

            # Process pages in batches to manage memory
            all_fonts = set()
            processed_pages = 0
            
            for batch_start in range(0, page_count, batch_size):
                batch_end = min(batch_start + batch_size, page_count)
                print(f"  📄 Processing pages {batch_start + 1}-{batch_end}...")
                
                # Process batch with parallel execution
                def process_page(page_num):
                    try:
                        page = doc[page_num]
                        result = self.extract_page_data(page, page_num + 1, pdf_path)
                        page.clean_contents()  # Clean page resources
                        return result
                    except Exception as e:
                        print(f"⚠️  Error in page {page_num + 1}: {str(e)}")
                        return None

                with concurrent.futures.ThreadPoolExecutor(max_workers=self.max_workers) as executor:
                    batch_pages = list(range(batch_start, batch_end))
                    futures = [executor.submit(process_page, p) for p in batch_pages]
                    
                    for future in concurrent.futures.as_completed(futures):
                        page_data = future.result()
                        if page_data:
                            # Extract essential info only for large files
                            if file_size_mb > 100:
                                # Keep only essential data for large files
                                essential_page_data = {
                                    'page_number': page_data['page_number'],
                                    'text': page_data['text'],
                                    'word_count': page_data['word_count'],
                                    'char_count': page_data['char_count'],
                                    'images': len(page_data['images']),  # Just count, not full data
                                    'tables': len(page_data['tables'])   # Just count, not full data
                                }
                                pdf_data['pages'].append(essential_page_data)
                            else:
                                pdf_data['pages'].append(page_data)
                            
                            # Accumulate totals
                            pdf_data['total_images'] += len(page_data['images'])
                            pdf_data['total_tables'] += len(page_data['tables'])
                            pdf_data['total_words'] += page_data['word_count']
                            pdf_data['total_characters'] += page_data['char_count']
                            all_fonts.update(page_data.get('fonts', []))
                            
                            processed_pages += 1
                
                # Force garbage collection after each batch
                gc.collect()
                
                # Check memory usage
                memory_usage = self.check_memory_usage()
                print(f"  💾 Memory usage: {memory_usage:.1f} MB")
                
                if self.is_memory_limit_exceeded():
                    print(f"  ⚠️  Memory limit ({self.memory_limit_mb} MB) exceeded, switching to minimal mode")
                    self.extract_formatting = False
                    # Process remaining pages in smaller batches
                    batch_size = max(1, batch_size // 2)

            pdf_data['fonts_used'] = list(all_fonts)

            # Extract tables using selective methods for large files
            print("  📊 Extracting tables...")
            if file_size_mb > 100:
                print("    🔧 Using optimized table extraction for large file...")
                pdf_data['extracted_tables'] = self.extract_tables_optimized(pdf_path, max_pages=min(50, page_count))
            else:
                pdf_data['extracted_tables'] = self.extract_tables_multiple_methods(pdf_path)

            doc.close()
            gc.collect()  # Final cleanup
            
            print(f"  ✅ PDF processed: {processed_pages} pages, {pdf_data['total_words']} words")
            return pdf_data

        except Exception as e:
            print(f"❌ Error processing PDF {pdf_path}: {str(e)}")
            return None

    def extract_page_data(self, page, page_num, pdf_path):
        """Extract data from one PDF page with memory optimization"""
        page_data = {
            'page_number': page_num,
            'text': '',
            'images': [],
            'tables': [],
            'fonts': [],
            'word_count': 0,
            'char_count': 0,
            'page_size': {
                'width': round(page.rect.width, 2),
                'height': round(page.rect.height, 2)
            },
            'rotation': page.rotation
        }
        
        # Only include detailed formatting for smaller files
        if self.extract_formatting:
            page_data['formatted_text'] = []
        
        try:
            # Extract text with optional formatting
            if self.extract_formatting:
                text_dict = page.get_text("dict")
                page_text = ""
                fonts_on_page = []

                for block in text_dict.get("blocks", []):
                    if "lines" in block:
                        for line in block["lines"]:
                            for span in line["spans"]:
                                text = span.get("text", "")
                                page_text += text
                                # Store formatted text only if enabled
                                page_data['formatted_text'].append({
                                    'text': text,
                                    'font': span.get('font', 'Unknown'),
                                    'size': round(span.get('size', 0), 2),
                                    'flags': span.get('flags', 0),
                                    'color': span.get('color', 0),
                                    'bbox': span.get('bbox', [0, 0, 0, 0])
                                })
                                fonts_on_page.append(span.get('font', 'Unknown'))
            else:
                # Faster text extraction without formatting
                page_text = page.get_text()
                fonts_on_page = []
            
            page_data['text'] = page_text
            page_data['fonts'] = list(set(fonts_on_page))
            page_data['word_count'] = len(page_text.split()) if page_text else 0
            page_data['char_count'] = len(page_text)

            # Extract images with memory management
            page_data['images'] = self.extract_images_from_page_optimized(page, page_num, pdf_path)

            # Extract tables
            page_data['tables'] = self.extract_tables_from_page(page)

        except Exception as e:
            print(f"⚠️  Error extracting page {page_num}: {str(e)}")
        
        return page_data

    def extract_images_from_page_optimized(self, page, page_num, pdf_path):
        """Extract images from a page with memory optimization"""
        images = []
        try:
            image_list = page.get_images()
            
            # Limit number of images processed for large files
            max_images = 10 if self.check_memory_usage() > 500 else len(image_list)
            
            for img_index, img in enumerate(image_list[:max_images]):
                try:
                    xref = img[0]
                    base_image = page.parent.extract_image(xref)
                    image_bytes = base_image["image"]
                    
                    # Skip very large images to save memory
                    if len(image_bytes) > 5 * 1024 * 1024:  # Skip images > 5MB
                        print(f"    ⚠️  Skipping large image ({len(image_bytes)//1024//1024} MB) on page {page_num}")
                        continue
                    
                    filename_base = os.path.splitext(os.path.basename(pdf_path))[0]
                    img_name = f"extracted_images/{filename_base}_page_{page_num}_img_{img_index+1}.png"
                    
                    # Save image to disk and free memory immediately
                    with open(img_name, "wb") as img_file:
                        img_file.write(image_bytes)
                    
                    img_info = {
                        'filename': img_name,
                        'width': base_image["width"],
                        'height': base_image["height"],
                        'colorspace': base_image.get("colorspace", "Unknown"),
                        'size_bytes': len(image_bytes)
                    }
                    images.append(img_info)
                    
                    # Clear image data from memory immediately
                    del image_bytes
                    del base_image
                    
                except Exception as e:
                    print(f"⚠️  Error extracting image {img_index} on page {page_num}: {str(e)}")
                    continue
                    
            # Force garbage collection for images
            gc.collect()
            
        except Exception as e:
            print(f"⚠️  No images on page {page_num}: {str(e)}")
        
        return images

    def is_likely_table(self, text_lines):
        """Heuristic to detect tables based on lines"""
        if len(text_lines) < 2:
            return False
        # Basic heuristic: presence of separators, formatting, numeric columns...
        separator_patterns = [r'\t', r'\s{2,}', r'\|', r',', r';']
        for pattern in separator_patterns:
            counts = [len(re.findall(pattern, line)) for line in text_lines]
            if len(set(counts)) == 1 and counts[0] > 0:
                return True
        # Numeric content heuristic
        numeric_lines = sum(1 for line in text_lines if re.search(r'\d+', line))
        if numeric_lines / len(text_lines) > 0.5:
            return True
        # Similar line lengths
        lengths = [len(line.strip()) for line in text_lines]
        if lengths:
            avg = sum(lengths)/len(lengths)
            similar = sum(1 for l in lengths if abs(l - avg) < avg*0.3)
            if similar / len(lengths) > 0.7:
                return True
        return False

    def extract_tables_from_page(self, page):
        """Detect tables heuristically from page"""
        tables = []
        try:
            text_dict = page.get_text("dict")
            for block in text_dict.get("blocks", []):
                if "lines" in block:
                    lines = block["lines"]
                    if len(lines) >= 3:
                        text_lines = []
                        for line in lines:
                            row_text = " ".join(span.get("text", "") for span in line["spans"])
                            text_lines.append(row_text.strip())
                        if self.is_likely_table(text_lines):
                            tables.append({
                                'bbox': block.get('bbox', [0, 0, 0, 0]),
                                'data': text_lines,
                                'method': 'heuristic',
                                'confidence': 'medium'
                            })
        except:
            pass
        return tables

    def extract_tables_multiple_methods(self, pdf_path):
        """Extract tables with multiple methods; can be expanded"""
        all_tables = []

        # Using tabula
        if HAS_TABULA:
            try:
                print("    📊 Trying tabula...")
                tables = tabula.read_pdf(pdf_path, pages='all', multiple_tables=True, silent=True)
                for i, table in enumerate(tables):
                    if not table.empty and table.shape[0] > 1 and table.shape[1] > 1:
                        if
Thought for 1s
Here's the complete corrected code with optimizations for handling large files:

📁 Complete Code Files
1. pdfextract.py - Main Extraction Engine
import concurrent.futures
import os
import json
import pandas as pd
from PIL import Image
import pytesseract
from collections import Counter
import zipfile
import tempfile
import warnings
import re
import fitz  # PyMuPDF
import gc  # Add garbage collection
import psutil  # Add memory monitoring

# Optional libraries
try:
    import pdfplumber
    HAS_PDFPLUMBER = True
except ImportError:
    HAS_PDFPLUMBER = False

try:
    import tabula
    HAS_TABULA = True
except ImportError:
    HAS_TABULA = False

try:
    import camelot
    HAS_CAMELOT = True
except ImportError:
    HAS_CAMELOT = False

try:
    import cv2
    import numpy as np
    HAS_CV2 = True
except ImportError:
    HAS_CV2 = False

try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import openpyxl
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False


class ComprehensiveDocumentExtractor:
    def __init__(self, max_workers=None, memory_limit_mb=1024, extract_formatting=False):
        self.extracted_data = {}
        self.supported_formats = ['.pdf', '.docx', '.pptx', '.xlsx', '.xls']
        self.memory_limit_mb = memory_limit_mb
        self.extract_formatting = extract_formatting  # Option to disable detailed formatting for large files
        self.max_workers = max_workers or min(8, os.cpu_count() or 4)  # Adaptive worker count
        self.setup_directories()

    def setup_directories(self):
        """Create necessary directories"""
        os.makedirs('extracted_images', exist_ok=True)
        os.makedirs('extracted_tables', exist_ok=True)

    def check_memory_usage(self):
        """Check current memory usage"""
        process = psutil.Process(os.getpid())
        memory_mb = process.memory_info().rss / 1024 / 1024
        return memory_mb

    def is_memory_limit_exceeded(self):
        """Check if memory limit is exceeded"""
        return self.check_memory_usage() > self.memory_limit_mb

    def extract_pdf_comprehensive(self, pdf_path):
        """Extract comprehensive data from PDF with optimized memory usage"""
        try:
            print(f"📄 Processing PDF: {os.path.basename(pdf_path)}")
            doc = fitz.open(pdf_path)
            
            # Check file size and adjust processing accordingly
            file_size_mb = os.path.getsize(pdf_path) / (1024*1024)
            page_count = len(doc)
            print(f"  📊 File size: {file_size_mb:.2f} MB, Pages: {page_count}")
            
            # Adjust processing based on file size
            if file_size_mb > 100:  # Large file optimizations
                print("  🔧 Large file detected - applying optimizations...")
                self.extract_formatting = False  # Disable detailed formatting
                batch_size = max(1, min(10, 50 // int(file_size_mb / 10)))  # Adaptive batch size
            else:
                batch_size = page_count
            
            pdf_data = {
                'filename': os.path.basename(pdf_path),
                'file_type': 'PDF',
                'metadata': {},
                'pages': [],
                'total_images': 0,
                'total_tables': 0,
                'fonts_used': [],
                'total_words': 0,
                'total_characters': 0,
                'page_count': page_count,
                'file_size_mb': round(file_size_mb, 2),
                'processing_mode': 'optimized' if file_size_mb > 100 else 'standard'
            }

            # Extract metadata safely
            try:
                metadata = doc.metadata
                if metadata:
                    pdf_data['metadata'] = {k: v for k, v in metadata.items() if v}
            except:
                pdf_data['metadata'] = {}

            # Process pages in batches to manage memory
            all_fonts = set()
            processed_pages = 0
            
            for batch_start in range(0, page_count, batch_size):
                batch_end = min(batch_start + batch_size, page_count)
                print(f"  📄 Processing pages {batch_start + 1}-{batch_end}...")
                
                # Process batch with parallel execution
                def process_page(page_num):
                    try:
                        page = doc[page_num]
                        result = self.extract_page_data(page, page_num + 1, pdf_path)
                        page.clean_contents()  # Clean page resources
                        return result
                    except Exception as e:
                        print(f"⚠️  Error in page {page_num + 1}: {str(e)}")
                        return None

                with concurrent.futures.ThreadPoolExecutor(max_workers=self.max_workers) as executor:
                    batch_pages = list(range(batch_start, batch_end))
                    futures = [executor.submit(process_page, p) for p in batch_pages]
                    
                    for future in concurrent.futures.as_completed(futures):
                        page_data = future.result()
                        if page_data:
                            # Extract essential info only for large files
                            if file_size_mb > 100:
                                # Keep only essential data for large files
                                essential_page_data = {
                                    'page_number': page_data['page_number'],
                                    'text': page_data['text'],
                                    'word_count': page_data['word_count'],
                                    'char_count': page_data['char_count'],
                                    'images': len(page_data['images']),  # Just count, not full data
                                    'tables': len(page_data['tables'])   # Just count, not full data
                                }
                                pdf_data['pages'].append(essential_page_data)
                            else:
                                pdf_data['pages'].append(page_data)
                            
                            # Accumulate totals
                            pdf_data['total_images'] += len(page_data['images'])
                            pdf_data['total_tables'] += len(page_data['tables'])
                            pdf_data['total_words'] += page_data['word_count']
                            pdf_data['total_characters'] += page_data['char_count']
                            all_fonts.update(page_data.get('fonts', []))
                            
                            processed_pages += 1
                
                # Force garbage collection after each batch
                gc.collect()
                
                # Check memory usage
                memory_usage = self.check_memory_usage()
                print(f"  💾 Memory usage: {memory_usage:.1f} MB")
                
                if self.is_memory_limit_exceeded():
                    print(f"  ⚠️  Memory limit ({self.memory_limit_mb} MB) exceeded, switching to minimal mode")
                    self.extract_formatting = False
                    # Process remaining pages in smaller batches
                    batch_size = max(1, batch_size // 2)

            pdf_data['fonts_used'] = list(all_fonts)

            # Extract tables using selective methods for large files
            print("  📊 Extracting tables...")
            if file_size_mb > 100:
                print("    🔧 Using optimized table extraction for large file...")
                pdf_data['extracted_tables'] = self.extract_tables_optimized(pdf_path, max_pages=min(50, page_count))
            else:
                pdf_data['extracted_tables'] = self.extract_tables_multiple_methods(pdf_path)

            doc.close()
            gc.collect()  # Final cleanup
            
            print(f"  ✅ PDF processed: {processed_pages} pages, {pdf_data['total_words']} words")
            return pdf_data

        except Exception as e:
            print(f"❌ Error processing PDF {pdf_path}: {str(e)}")
            return None

    def extract_page_data(self, page, page_num, pdf_path):
        """Extract data from one PDF page with memory optimization"""
        page_data = {
            'page_number': page_num,
            'text': '',
            'images': [],
            'tables': [],
            'fonts': [],
            'word_count': 0,
            'char_count': 0,
            'page_size': {
                'width': round(page.rect.width, 2),
                'height': round(page.rect.height, 2)
            },
            'rotation': page.rotation
        }
        
        # Only include detailed formatting for smaller files
        if self.extract_formatting:
            page_data['formatted_text'] = []
        
        try:
            # Extract text with optional formatting
            if self.extract_formatting:
                text_dict = page.get_text("dict")
                page_text = ""
                fonts_on_page = []

                for block in text_dict.get("blocks", []):
                    if "lines" in block:
                        for line in block["lines"]:
                            for span in line["spans"]:
                                text = span.get("text", "")
                                page_text += text
                                # Store formatted text only if enabled
                                page_data['formatted_text'].append({
                                    'text': text,
                                    'font': span.get('font', 'Unknown'),
                                    'size': round(span.get('size', 0), 2),
                                    'flags': span.get('flags', 0),
                                    'color': span.get('color', 0),
                                    'bbox': span.get('bbox', [0, 0, 0, 0])
                                })
                                fonts_on_page.append(span.get('font', 'Unknown'))
            else:
                # Faster text extraction without formatting
                page_text = page.get_text()
                fonts_on_page = []
            
            page_data['text'] = page_text
            page_data['fonts'] = list(set(fonts_on_page))
            page_data['word_count'] = len(page_text.split()) if page_text else 0
            page_data['char_count'] = len(page_text)

            # Extract images with memory management
            page_data['images'] = self.extract_images_from_page_optimized(page, page_num, pdf_path)

            # Extract tables
            page_data['tables'] = self.extract_tables_from_page(page)

        except Exception as e:
            print(f"⚠️  Error extracting page {page_num}: {str(e)}")
        
        return page_data

    def extract_images_from_page_optimized(self, page, page_num, pdf_path):
        """Extract images from a page with memory optimization"""
        images = []
        try:
            image_list = page.get_images()
            
            # Limit number of images processed for large files
            max_images = 10 if self.check_memory_usage() > 500 else len(image_list)
            
            for img_index, img in enumerate(image_list[:max_images]):
                try:
                    xref = img[0]
                    base_image = page.parent.extract_image(xref)
                    image_bytes = base_image["image"]
                    
                    # Skip very large images to save memory
                    if len(image_bytes) > 5 * 1024 * 1024:  # Skip images > 5MB
                        print(f"    ⚠️  Skipping large image ({len(image_bytes)//1024//1024} MB) on page {page_num}")
                        continue
                    
                    filename_base = os.path.splitext(os.path.basename(pdf_path))[0]
                    img_name = f"extracted_images/{filename_base}_page_{page_num}_img_{img_index+1}.png"
                    
                    # Save image to disk and free memory immediately
                    with open(img_name, "wb") as img_file:
                        img_file.write(image_bytes)
                    
                    img_info = {
                        'filename': img_name,
                        'width': base_image["width"],
                        'height': base_image["height"],
                        'colorspace': base_image.get("colorspace", "Unknown"),
                        'size_bytes': len(image_bytes)
                    }
                    images.append(img_info)
                    
                    # Clear image data from memory immediately
                    del image_bytes
                    del base_image
                    
                except Exception as e:
                    print(f"⚠️  Error extracting image {img_index} on page {page_num}: {str(e)}")
                    continue
                    
            # Force garbage collection for images
            gc.collect()
            
        except Exception as e:
            print(f"⚠️  No images on page {page_num}: {str(e)}")
        
        return images

    def is_likely_table(self, text_lines):
        """Heuristic to detect tables based on lines"""
        if len(text_lines) < 2:
            return False
        # Basic heuristic: presence of separators, formatting, numeric columns...
        separator_patterns = [r'\t', r'\s{2,}', r'\|', r',', r';']
        for pattern in separator_patterns:
            counts = [len(re.findall(pattern, line)) for line in text_lines]
            if len(set(counts)) == 1 and counts[0] > 0:
                return True
        # Numeric content heuristic
        numeric_lines = sum(1 for line in text_lines if re.search(r'\d+', line))
        if numeric_lines / len(text_lines) > 0.5:
            return True
        # Similar line lengths
        lengths = [len(line.strip()) for line in text_lines]
        if lengths:
            avg = sum(lengths)/len(lengths)
            similar = sum(1 for l in lengths if abs(l - avg) < avg*0.3)
            if similar / len(lengths) > 0.7:
                return True
        return False

    def extract_tables_from_page(self, page):
        """Detect tables heuristically from page"""
        tables = []
        try:
            text_dict = page.get_text("dict")
            for block in text_dict.get("blocks", []):
                if "lines" in block:
                    lines = block["lines"]
                    if len(lines) >= 3:
                        text_lines = []
                        for line in lines:
                            row_text = " ".join(span.get("text", "") for span in line["spans"])
                            text_lines.append(row_text.strip())
                        if self.is_likely_table(text_lines):
                            tables.append({
                                'bbox': block.get('bbox', [0, 0, 0, 0]),
                                'data': text_lines,
                                'method': 'heuristic',
                                'confidence': 'medium'
                            })
        except:
            pass
        return tables

    def extract_tables_multiple_methods(self, pdf_path):
        """Extract tables with multiple methods; can be expanded"""
        all_tables = []

        # Using tabula
        if HAS_TABULA:
            try:
                print("    📊 Trying tabula...")
                tables = tabula.read_pdf(pdf_path, pages='all', multiple_tables=True, silent=True)
                for i, table in enumerate(tables):
                    if not table.empty and table.shape[0] > 1 and table.shape[1] > 1:
                        if self.validate_extracted_table(table):
                            all_tables.append({
                                'method': 'tabula',
                                'table_index': i,
                                'data': table.to_dict('records'),
                                'shape': table.shape,
                                'columns': table.columns.tolist(),
                                'confidence': 'high'
                            })
            except Exception as e:
                print(f"    ⚠️  Tabula error: {str(e)}")
        # Using camelot
        if HAS_CAMELOT:
            try:
                print("    📊 Trying camelot...")
                tables = camelot.read_pdf(pdf_path, pages='all', flavor='lattice')
                for i, table in enumerate(tables):
                    if not table.df.empty and table.df.shape[0] > 1 and table.df.shape[1] > 1:
                        accuracy = getattr(table, 'accuracy', 0)
                        if accuracy > 50:
                            all_tables.append({
                                'method': 'camelot_lattice',
                                'table_index': i,
                                'data': table.df.to_dict('records'),
                                'shape': table.df.shape,
                                'accuracy': accuracy,
                                'confidence': 'high' if accuracy > 80 else 'medium'
                            })
            except Exception as e:
                print(f"    ⚠️  Camelot error: {str(e)}")
        # Using pdfplumber
        if HAS_PDFPLUMBER:
            try:
                print("    📊 Trying pdfplumber...")
                with pdfplumber.open(pdf_path) as pdf:
                    for p_idx, page in enumerate(pdf.pages):
                        tables = page.extract_tables()
                        for t_idx, table in enumerate(tables):
                            if table and len(table) > 1 and len(table[0]) > 1:
                                if self.validate_pdfplumber_table(table):
                                    all_tables.append({
                                        'method': 'pdfplumber',
                                        'page': p_idx + 1,
                                        'table_index': t_idx,
                                        'data': table,
                                        'shape': (len(table), len(table[0]) if table else 0),
                                        'confidence': 'medium'
                                    })
            except:
                pass
        return all_tables

    def validate_extracted_table(self, df):
        """Validate pandas DataFrame as table"""
        if df.empty or df.shape[0] < 2 or df.shape[1] < 2:
            return False
        # Must have content
        non_empty = df.astype(str).values.flatten()
        non_empty_cells = [cell for cell in non_empty if cell.strip() and cell.strip() != 'nan']
        if len(non_empty_cells) < df.size * 0.3:
            return False
        if len(set(non_empty_cells)) < 2:
            return False
        return True

    def validate_pdfplumber_table(self, table):
        """Validate pdfplumber table"""
        if not table or len(table) < 2:
            return False
        total_cells = 0
        non_empty_cells = 0
        unique_vals = set()
        for row in table:
            for cell in row:
                total_cells += 1
                if cell and cell.strip():
                    non_empty_cells += 1
                    unique_vals.add(cell.strip())
        if non_empty_cells < total_cells * 0.3:
            return False
        if len(unique_vals) < 2:
            return False
        return True
    
    def extract_tables_optimized(self, pdf_path, max_pages=50):
        """Optimized table extraction for large files"""
        all_tables = []
        
        # For large files, use only the most reliable method and limit pages
        print(f"    📊 Optimized extraction - processing first {max_pages} pages only...")
        
        # Try pdfplumber first (most memory efficient)
        if HAS_PDFPLUMBER:
            try:
                with pdfplumber.open(pdf_path) as pdf:
                    pages_to_process = min(max_pages, len(pdf.pages))
                    for p_idx in range(pages_to_process):
                        try:
                            page = pdf.pages[p_idx]
                            tables = page.extract_tables()
                            for t_idx, table in enumerate(tables):
                                if table and len(table) > 1 and len(table[0]) > 1:
                                    if self.validate_pdfplumber_table(table):
                                        # Limit table size for memory efficiency
                                        if len(table) > 100:
                                            table = table[:100]  # Keep first 100 rows
                                        all_tables.append({
                                            'method': 'pdfplumber_optimized',
                                            'page': p_idx + 1,
                                            'table_index': t_idx,
                                            'data': table,
                                            'shape': (len(table), len(table[0]) if table else 0),
                                            'confidence': 'medium',
                                            'note': 'Limited for large file processing'
                                        })
                        except Exception as e:
                            print(f"    ⚠️  Error on page {p_idx + 1}: {str(e)}")
                            continue
                        
                        # Check memory usage periodically
                        if p_idx % 10 == 0:
                            gc.collect()
                            if self.is_memory_limit_exceeded():
                                print(f"    ⚠️  Memory limit reached at page {p_idx + 1}")
                                break
            except Exception as e:
                print(f"    ⚠️  Pdfplumber optimized error: {str(e)}")
        
        # If no tables found with pdfplumber, try tabula on first few pages only
        if not all_tables and HAS_TABULA:
            try:
                print("    📊 Trying tabula on first 10 pages...")
                pages_str = '1-10'  # Limit to first 10 pages
                tables = tabula.read_pdf(pdf_path, pages=pages_str, multiple_tables=True, silent=True)
                for i, table in enumerate(tables):
                    if not table.empty and table.shape[0] > 1 and table.shape[1] > 1:
                        if self.validate_extracted_table(table):
                            # Limit table size
                            if table.shape[0] > 100:
                                table = table.head(100)
                            all_tables.append({
                                'method': 'tabula_optimized',
                                'table_index': i,
                                'data': table.to_dict('records'),
                                'shape': table.shape,
                                'columns': table.columns.tolist(),
                                'confidence': 'high',
                                'note': 'Limited to first 10 pages and 100 rows'
                            })
            except Exception as e:
                print(f"    ⚠️  Tabula optimized error: {str(e)}")
        
        return all_tables
    
    def extract_docx(self, docx_path):
        """Extract data from DOCX files"""
        if not HAS_DOCX:
            print("❌ python-docx not available. Skipping DOCX processing.")
            return None
            
        try:
            print(f"📄 Processing DOCX: {os.path.basename(docx_path)}")
            doc = Document(docx_path)
            
            docx_data = {
                'filename': os.path.basename(docx_path),
                'file_type': 'DOCX',
                'paragraphs': [],
                'tables': [],
                'total_words': 0,
                'total_characters': 0,
                'file_size_mb': round(os.path.getsize(docx_path) / (1024*1024), 2)
            }
            
            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    docx_data['paragraphs'].append({
                        'text': para.text,
                        'style': para.style.name if para.style else 'Normal'
                    })
                    docx_data['total_words'] += len(para.text.split())
                    docx_data['total_characters'] += len(para.text)
            
            # Extract tables
            for table_idx, table in enumerate(doc.tables):
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                
                if table_data and len(table_data) > 1:
                    docx_data['tables'].append({
                        'table_index': table_idx,
                        'data': table_data,
                        'shape': (len(table_data), len(table_data[0]) if table_data else 0)
                    })
            
            print(f"  ✅ DOCX processed: {len(docx_data['paragraphs'])} paragraphs, {docx_data['total_words']} words")
            return docx_data
            
        except Exception as e:
            print(f"❌ Error processing DOCX {docx_path}: {str(e)}")
            return None
    
    def extract_pptx(self, pptx_path):
        """Extract data from PPTX files"""
        if not HAS_PPTX:
            print("❌ python-pptx not available. Skipping PPTX processing.")
            return None
            
        try:
            print(f"📄 Processing PPTX: {os.path.basename(pptx_path)}")
            prs = Presentation(pptx_path)
            
            pptx_data = {
                'filename': os.path.basename(pptx_path),
                'file_type': 'PPTX',
                'slides': [],
                'total_slides': len(prs.slides),
                'total_words': 0,
                'total_characters': 0,
                'file_size_mb': round(os.path.getsize(pptx_path) / (1024*1024), 2)
            }
            
            for slide_idx, slide in enumerate(prs.slides):
                slide_data = {
                    'slide_number': slide_idx + 1,
                    'text_content': [],
                    'tables': []
                }
                
                # Extract text from shapes
                for shape in slide.shapes:
                    try:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_data['text_content'].append(shape.text.strip())
                            pptx_data['total_words'] += len(shape.text.split())
                            pptx_data['total_characters'] += len(shape.text)
                        
                        # Extract tables
                        if hasattr(shape, "table") and shape.table:
                            table_data = []
                            for row in shape.table.rows:
                                row_data = [cell.text.strip() for cell in row.cells]
                                table_data.append(row_data)
                            if table_data and len(table_data) > 1:
                                slide_data['tables'].append(table_data)
                    except Exception as e:
                        continue
                
                pptx_data['slides'].append(slide_data)
            
            print(f"  ✅ PPTX processed: {pptx_data['total_slides']} slides, {pptx_data['total_words']} words")
            return pptx_data
            
        except Exception as e:
            print(f"❌ Error processing PPTX {pptx_path}: {str(e)}")
            return None
    
    def extract_excel(self, excel_path):
        """Extract data from Excel files with large file optimization"""
        try:
            print(f"📄 Processing Excel: {os.path.basename(excel_path)}")
            
            file_size_mb = os.path.getsize(excel_path) / (1024*1024)
            print(f"  📊 File size: {file_size_mb:.2f} MB")
            
            excel_data = {
                'filename': os.path.basename(excel_path),
                'file_type': 'Excel',
                'sheets': [],
                'file_size_mb': round(file_size_mb, 2),
                'processing_mode': 'optimized' if file_size_mb > 50 else 'standard'
            }
            
            # For large files, implement streaming and limits
            max_rows_per_sheet = 1000 if file_size_mb > 50 else None
            chunk_size = 100 if file_size_mb > 100 else None
            
            # Get sheet names without loading data
            try:
                xl_file = pd.ExcelFile(excel_path)
                sheet_names = xl_file.sheet_names
                
                for sheet_name in sheet_names:
                    try:
                        print(f"    📋 Processing sheet: {sheet_name}")
                        
                        # For large files, read in chunks or with row limits
                        if chunk_size and max_rows_per_sheet:
                            # Read with both chunking and row limits for very large files
                            print(f"    🔧 Using chunked reading with {chunk_size} rows per chunk, max {max_rows_per_sheet} rows")
                            
                            all_chunks = []
                            rows_read = 0
                            
                            try:
                                for chunk in pd.read_excel(excel_path, sheet_name=sheet_name, chunksize=chunk_size):
                                    all_chunks.append(chunk)
                                    rows_read += len(chunk)
                                    
                                    if rows_read >= max_rows_per_sheet:
                                        break
                                    
                                    # Check memory usage periodically
                                    if len(all_chunks) % 5 == 0:
                                        gc.collect()
                                        if self.is_memory_limit_exceeded():
                                            print(f"    ⚠️  Memory limit reached, stopping at {rows_read} rows")
                                            break
                                
                                # Combine chunks
                                if all_chunks:
                                    df = pd.concat(all_chunks, ignore_index=True)
                                else:
                                    df = pd.DataFrame()
                                    
                            except Exception as e:
                                print(f"    ⚠️  Chunked reading failed, trying row limit: {str(e)}")
                                # Fallback to simple row limit
                                df = pd.read_excel(excel_path, sheet_name=sheet_name, nrows=max_rows_per_sheet)
                                
                        elif max_rows_per_sheet:
                            # Read with row limit only
                            print(f"    🔧 Using row limit: {max_rows_per_sheet} rows")
                            df = pd.read_excel(excel_path, sheet_name=sheet_name, nrows=max_rows_per_sheet)
                            
                        else:
                            # Standard reading for smaller files
                            df = pd.read_excel(excel_path, sheet_name=sheet_name)
                        
                        # Process the dataframe
                        if not df.empty:
                            # For large sheets, limit the data we store
                            sample_rows = min(100, len(df)) if file_size_mb > 50 else len(df)
                            
                            sheet_data = {
                                'sheet_name': sheet_name,
                                'shape': df.shape,
                                'columns': df.columns.tolist(),
                                'data': df.head(sample_rows).to_dict('records'),
                                'data_types': df.dtypes.astype(str).to_dict()
                            }
                            
                            # Add summary statistics only for smaller files or limited rows
                            try:
                                if file_size_mb <= 50 or df.shape[0] <= 1000:
                                    sheet_data['summary'] = df.describe(include='all').to_dict()
                                else:
                                    # Limited summary for large files
                                    numeric_cols = df.select_dtypes(include=['number']).columns
                                    if len(numeric_cols) > 0:
                                        sheet_data['summary'] = df[numeric_cols].describe().to_dict()
                                    else:
                                        sheet_data['summary'] = {'note': 'Summary skipped for large non-numeric data'}
                            except Exception as e:
                                sheet_data['summary'] = {'error': f'Summary calculation failed: {str(e)}'}
                            
                            # Add processing notes for large files
                            if max_rows_per_sheet and df.shape[0] >= max_rows_per_sheet:
                                sheet_data['processing_note'] = f'Limited to first {max_rows_per_sheet} rows due to file size'
                            if file_size_mb > 50:
                                sheet_data['sample_note'] = f'Data sample limited to first {sample_rows} rows for display'
                            
                            excel_data['sheets'].append(sheet_data)
                            
                            # Clear dataframe from memory
                            del df
                            gc.collect()
                        
                        else:
                            # Empty sheet
                            sheet_data = {
                                'sheet_name': sheet_name,
                                'shape': (0, 0),
                                'columns': [],
                                'data': [],
                                'summary': {},
                                'data_types': {},
                                'processing_note': 'Sheet is empty'
                            }
                            excel_data['sheets'].append(sheet_data)
                        
                    except Exception as e:
                        print(f"    ⚠️  Error reading sheet '{sheet_name}': {str(e)}")
                        # Add error sheet info
                        error_sheet = {
                            'sheet_name': sheet_name,
                            'shape': (0, 0),
                            'columns': [],
                            'data': [],
                            'summary': {},
                            'data_types': {},
                            'processing_note': f'Error reading sheet: {str(e)}'
                        }
                        excel_data['sheets'].append(error_sheet)
                        continue
                
                # Final cleanup
                xl_file.close()
                gc.collect()
                
            except Exception as e:
                print(f"    ❌ Error opening Excel file: {str(e)}")
                return None
            
            print(f"  ✅ Excel processed: {len(excel_data['sheets'])} sheets")
            return excel_data
            
        except Exception as e:
            print(f"❌ Error processing Excel {excel_path}: {str(e)}")
            return None
    
    def process_files(self, file_paths):
        """Process all uploaded files"""
        results = {}
        
        print(f"\n🔄 Processing {len(file_paths)} file(s)...")
        print("=" * 60)
        
        for i, file_path in enumerate(file_paths, 1):
            file_ext = os.path.splitext(file_path)[1].lower()
            
            print(f"\n[{i}/{len(file_paths)}] Processing: {file_path}")
            
            try:
                if file_ext == '.pdf':
                    result = self.extract_pdf_comprehensive(file_path)
                elif file_ext == '.docx':
                    result = self.extract_docx(file_path)
                elif file_ext == '.pptx':
                    result = self.extract_pptx(file_path)
                elif file_ext in ['.xlsx', '.xls']:
                    result = self.extract_excel(file_path)
                else:
                    print(f"❌ Unsupported file format: {file_ext}")
                    continue
                
                if result:
                    results[file_path] = result
                    print(f"✅ Successfully processed: {file_path}")
                else:
                    print(f"❌ Failed to process: {file_path}")
                    
            except Exception as e:
                print(f"❌ Error processing {file_path}: {str(e)}")
                continue
        
        return results
    
    def save_results(self, results, output_format='json'):
        """Save extraction results with better formatting"""
        if not results:
            print("❌ No results to save.")
            return
        
        try:
            # Save as JSON
            with open('extraction_results.json', 'w', encoding='utf-8') as f:
                json.dump(results, f, indent=2, ensure_ascii=False, default=str)
            print("✅ Results saved as: extraction_results.json")
            
            # Create summary report
            self.create_summary_report(results)
            
            # Create CSV summaries for tables
            self.save_tables_as_csv(results)
            
            # Check if there are extracted images
            if os.path.exists('extracted_images') and os.listdir('extracted_images'):
                print("🖼️  Images extracted to 'extracted_images' folder")
                
        except Exception as e:
            print(f"❌ Error saving results: {str(e)}")
    
    def save_tables_as_csv(self, results):
        """Save extracted tables as CSV files"""
        table_count = 0
        
        for file_path, data in results.items():
            filename_base = os.path.splitext(data['filename'])[0]
            
            # Save tables from different extraction methods
            if 'extracted_tables' in data:
                for i, table in enumerate(data['extracted_tables']):
                    try:
                        if 'data' in table and table['data']:
                            csv_filename = f"extracted_tables/{filename_base}_table_{i+1}_{table['method']}.csv"
                            
                            if table['method'] in ['tabula', 'camelot_lattice']:
                                # These are already in records format
                                df = pd.DataFrame(table['data'])
                            else:
                                # Convert list of lists to DataFrame
                                df = pd.DataFrame(table['data'])
                            
                            df.to_csv(csv_filename, index=False)
                            
Here's the complete corrected code with optimizations for handling large files:

📁 Complete Code Files
1. pdfextract.py - Main Extraction Engine
import concurrent.futures
import os
import json
import pandas as pd
from PIL import Image
import pytesseract
from collections import Counter
import zipfile
import tempfile
import warnings
import re
import fitz  # PyMuPDF
import gc  # Add garbage collection
import psutil  # Add memory monitoring

# Optional libraries
try:
    import pdfplumber
    HAS_PDFPLUMBER = True
except ImportError:
    HAS_PDFPLUMBER = False

try:
    import tabula
    HAS_TABULA = True
except ImportError:
    HAS_TABULA = False

try:
    import camelot
    HAS_CAMELOT = True
except ImportError:
    HAS_CAMELOT = False

try:
    import cv2
    import numpy as np
    HAS_CV2 = True
except ImportError:
    HAS_CV2 = False

try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import openpyxl
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False


class ComprehensiveDocumentExtractor:
    def __init__(self, max_workers=None, memory_limit_mb=1024, extract_formatting=False):
        self.extracted_data = {}
        self.supported_formats = ['.pdf', '.docx', '.pptx', '.xlsx', '.xls']
        self.memory_limit_mb = memory_limit_mb
        self.extract_formatting = extract_formatting  # Option to disable detailed formatting for large files
        self.max_workers = max_workers or min(8, os.cpu_count() or 4)  # Adaptive worker count
        self.setup_directories()

    def setup_directories(self):
        """Create necessary directories"""
        os.makedirs('extracted_images', exist_ok=True)
        os.makedirs('extracted_tables', exist_ok=True)

    def check_memory_usage(self):
        """Check current memory usage"""
        process = psutil.Process(os.getpid())
        memory_mb = process.memory_info().rss / 1024 / 1024
        return memory_mb

    def is_memory_limit_exceeded(self):
        """Check if memory limit is exceeded"""
        return self.check_memory_usage() > self.memory_limit_mb

    def extract_pdf_comprehensive(self, pdf_path):
        """Extract comprehensive data from PDF with optimized memory usage"""
        try:
            print(f"📄 Processing PDF: {os.path.basename(pdf_path)}")
            doc = fitz.open(pdf_path)
            
            # Check file size and adjust processing accordingly
            file_size_mb = os.path.getsize(pdf_path) / (1024*1024)
            page_count = len(doc)
            print(f"  📊 File size: {file_size_mb:.2f} MB, Pages: {page_count}")
            
            # Adjust processing based on file size
            if file_size_mb > 100:  # Large file optimizations
                print("  🔧 Large file detected - applying optimizations...")
                self.extract_formatting = False  # Disable detailed formatting
                batch_size = max(1, min(10, 50 // int(file_size_mb / 10)))  # Adaptive batch size
            else:
                batch_size = page_count
            
            pdf_data = {
                'filename': os.path.basename(pdf_path),
                'file_type': 'PDF',
                'metadata': {},
                'pages': [],
                'total_images': 0,
                'total_tables': 0,
                'fonts_used': [],
                'total_words': 0,
                'total_characters': 0,
                'page_count': page_count,
                'file_size_mb': round(file_size_mb, 2),
                'processing_mode': 'optimized' if file_size_mb > 100 else 'standard'
            }

            # Extract metadata safely
            try:
                metadata = doc.metadata
                if metadata:
                    pdf_data['metadata'] = {k: v for k, v in metadata.items() if v}
            except:
                pdf_data['metadata'] = {}

            # Process pages in batches to manage memory
            all_fonts = set()
            processed_pages = 0
            
            for batch_start in range(0, page_count, batch_size):
                batch_end = min(batch_start + batch_size, page_count)
                print(f"  📄 Processing pages {batch_start + 1}-{batch_end}...")
                
                # Process batch with parallel execution
                def process_page(page_num):
                    try:
                        page = doc[page_num]
                        result = self.extract_page_data(page, page_num + 1, pdf_path)
                        page.clean_contents()  # Clean page resources
                        return result
                    except Exception as e:
                        print(f"⚠️  Error in page {page_num + 1}: {str(e)}")
                        return None

                with concurrent.futures.ThreadPoolExecutor(max_workers=self.max_workers) as executor:
                    batch_pages = list(range(batch_start, batch_end))
                    futures = [executor.submit(process_page, p) for p in batch_pages]
                    
                    for future in concurrent.futures.as_completed(futures):
                        page_data = future.result()
                        if page_data:
                            # Extract essential info only for large files
                            if file_size_mb > 100:
                                # Keep only essential data for large files
                                essential_page_data = {
                                    'page_number': page_data['page_number'],
                                    'text': page_data['text'],
                                    'word_count': page_data['word_count'],
                                    'char_count': page_data['char_count'],
                                    'images': len(page_data['images']),  # Just count, not full data
                                    'tables': len(page_data['tables'])   # Just count, not full data
                                }
                                pdf_data['pages'].append(essential_page_data)
                            else:
                                pdf_data['pages'].append(page_data)
                            
                            # Accumulate totals
                            pdf_data['total_images'] += len(page_data['images'])
                            pdf_data['total_tables'] += len(page_data['tables'])
                            pdf_data['total_words'] += page_data['word_count']
                            pdf_data['total_characters'] += page_data['char_count']
                            all_fonts.update(page_data.get('fonts', []))
                            
                            processed_pages += 1
                
                # Force garbage collection after each batch
                gc.collect()
                
                # Check memory usage
                memory_usage = self.check_memory_usage()
                print(f"  💾 Memory usage: {memory_usage:.1f} MB")
                
                if self.is_memory_limit_exceeded():
                    print(f"  ⚠️  Memory limit ({self.memory_limit_mb} MB) exceeded, switching to minimal mode")
                    self.extract_formatting = False
                    # Process remaining pages in smaller batches
                    batch_size = max(1, batch_size // 2)

            pdf_data['fonts_used'] = list(all_fonts)

            # Extract tables using selective methods for large files
            print("  📊 Extracting tables...")
            if file_size_mb > 100:
                print("    🔧 Using optimized table extraction for large file...")
                pdf_data['extracted_tables'] = self.extract_tables_optimized(pdf_path, max_pages=min(50, page_count))
            else:
                pdf_data['extracted_tables'] = self.extract_tables_multiple_methods(pdf_path)

            doc.close()
            gc.collect()  # Final cleanup
            
            print(f"  ✅ PDF processed: {processed_pages} pages, {pdf_data['total_words']} words")
            return pdf_data

        except Exception as e:
            print(f"❌ Error processing PDF {pdf_path}: {str(e)}")
            return None

    def extract_page_data(self, page, page_num, pdf_path):
        """Extract data from one PDF page with memory optimization"""
        page_data = {
            'page_number': page_num,
            'text': '',
            'images': [],
            'tables': [],
            'fonts': [],
            'word_count': 0,
            'char_count': 0,
            'page_size': {
                'width': round(page.rect.width, 2),
                'height': round(page.rect.height, 2)
            },
            'rotation': page.rotation
        }
        
        # Only include detailed formatting for smaller files
        if self.extract_formatting:
            page_data['formatted_text'] = []
        
        try:
            # Extract text with optional formatting
            if self.extract_formatting:
                text_dict = page.get_text("dict")
                page_text = ""
                fonts_on_page = []

                for block in text_dict.get("blocks", []):
                    if "lines" in block:
                        for line in block["lines"]:
                            for span in line["spans"]:
                                text = span.get("text", "")
                                page_text += text
                                # Store formatted text only if enabled
                                page_data['formatted_text'].append({
                                    'text': text,
                                    'font': span.get('font', 'Unknown'),
                                    'size': round(span.get('size', 0), 2),
                                    'flags': span.get('flags', 0),
                                    'color': span.get('color', 0),
                                    'bbox': span.get('bbox', [0, 0, 0, 0])
                                })
                                fonts_on_page.append(span.get('font', 'Unknown'))
            else:
                # Faster text extraction without formatting
                page_text = page.get_text()
                fonts_on_page = []
            
            page_data['text'] = page_text
            page_data['fonts'] = list(set(fonts_on_page))
            page_data['word_count'] = len(page_text.split()) if page_text else 0
            page_data['char_count'] = len(page_text)

            # Extract images with memory management
            page_data['images'] = self.extract_images_from_page_optimized(page, page_num, pdf_path)

            # Extract tables
            page_data['tables'] = self.extract_tables_from_page(page)

        except Exception as e:
            print(f"⚠️  Error extracting page {page_num}: {str(e)}")
        
        return page_data

    def extract_images_from_page_optimized(self, page, page_num, pdf_path):
        """Extract images from a page with memory optimization"""
        images = []
        try:
            image_list = page.get_images()
            
            # Limit number of images processed for large files
            max_images = 10 if self.check_memory_usage() > 500 else len(image_list)
            
            for img_index, img in enumerate(image_list[:max_images]):
                try:
                    xref = img[0]
                    base_image = page.parent.extract_image(xref)
                    image_bytes = base_image["image"]
                    
                    # Skip very large images to save memory
                    if len(image_bytes) > 5 * 1024 * 1024:  # Skip images > 5MB
                        print(f"    ⚠️  Skipping large image ({len(image_bytes)//1024//1024} MB) on page {page_num}")
                        continue
                    
                    filename_base = os.path.splitext(os.path.basename(pdf_path))[0]
                    img_name = f"extracted_images/{filename_base}_page_{page_num}_img_{img_index+1}.png"
                    
                    # Save image to disk and free memory immediately
                    with open(img_name, "wb") as img_file:
                        img_file.write(image_bytes)
                    
                    img_info = {
                        'filename': img_name,
                        'width': base_image["width"],
                        'height': base_image["height"],
                        'colorspace': base_image.get("colorspace", "Unknown"),
                        'size_bytes': len(image_bytes)
                    }
                    images.append(img_info)
                    
                    # Clear image data from memory immediately
                    del image_bytes
                    del base_image
                    
                except Exception as e:
                    print(f"⚠️  Error extracting image {img_index} on page {page_num}: {str(e)}")
                    continue
                    
            # Force garbage collection for images
            gc.collect()
            
        except Exception as e:
            print(f"⚠️  No images on page {page_num}: {str(e)}")
        
        return images

    def is_likely_table(self, text_lines):
        """Heuristic to detect tables based on lines"""
        if len(text_lines) < 2:
            return False
        # Basic heuristic: presence of separators, formatting, numeric columns...
        separator_patterns = [r'\t', r'\s{2,}', r'\|', r',', r';']
        for pattern in separator_patterns:
            counts = [len(re.findall(pattern, line)) for line in text_lines]
            if len(set(counts)) == 1 and counts[0] > 0:
                return True
        # Numeric content heuristic
        numeric_lines = sum(1 for line in text_lines if re.search(r'\d+', line))
        if numeric_lines / len(text_lines) > 0.5:
            return True
        # Similar line lengths
        lengths = [len(line.strip()) for line in text_lines]
        if lengths:
            avg = sum(lengths)/len(lengths)
            similar = sum(1 for l in lengths if abs(l - avg) < avg*0.3)
            if similar / len(lengths) > 0.7:
                return True
        return False

    def extract_tables_from_page(self, page):
        """Detect tables heuristically from page"""
        tables = []
        try:
            text_dict = page.get_text("dict")
            for block in text_dict.get("blocks", []):
                if "lines" in block:
                    lines = block["lines"]
                    if len(lines) >= 3:
                        text_lines = []
                        for line in lines:
                            row_text = " ".join(span.get("text", "") for span in line["spans"])
                            text_lines.append(row_text.strip())
                        if self.is_likely_table(text_lines):
                            tables.append({
                                'bbox': block.get('bbox', [0, 0, 0, 0]),
                                'data': text_lines,
                                'method': 'heuristic',
                                'confidence': 'medium'
                            })
        except:
            pass
        return tables

    def extract_tables_multiple_methods(self, pdf_path):
        """Extract tables with multiple methods; can be expanded"""
        all_tables = []

        # Using tabula
        if HAS_TABULA:
            try:
                print("    📊 Trying tabula...")
                tables = tabula.read_pdf(pdf_path, pages='all', multiple_tables=True, silent=True)
                for i, table in enumerate(tables):
                    if not table.empty and table.shape[0] > 1 and table.shape[1] > 1:
                        if self.validate_extracted_table(table):
                            all_tables.append({
                                'method': 'tabula',
                                'table_index': i,
                                'data': table.to_dict('records'),
                                'shape': table.shape,
                                'columns': table.columns.tolist(),
                                'confidence': 'high'
                            })
            except Exception as e:
                print(f"    ⚠️  Tabula error: {str(e)}")
        # Using camelot
        if HAS_CAMELOT:
            try:
                print("    📊 Trying camelot...")
                tables = camelot.read_pdf(pdf_path, pages='all', flavor='lattice')
                for i, table in enumerate(tables):
                    if not table.df.empty and table.df.shape[0] > 1 and table.df.shape[1] > 1:
                        accuracy = getattr(table, 'accuracy', 0)
                        if accuracy > 50:
                            all_tables.append({
                                'method': 'camelot_lattice',
                                'table_index': i,
                                'data': table.df.to_dict('records'),
                                'shape': table.df.shape,
                                'accuracy': accuracy,
                                'confidence': 'high' if accuracy > 80 else 'medium'
                            })
            except Exception as e:
                print(f"    ⚠️  Camelot error: {str(e)}")
        # Using pdfplumber
        if HAS_PDFPLUMBER:
            try:
                print("    📊 Trying pdfplumber...")
                with pdfplumber.open(pdf_path) as pdf:
                    for p_idx, page in enumerate(pdf.pages):
                        tables = page.extract_tables()
                        for t_idx, table in enumerate(tables):
                            if table and len(table) > 1 and len(table[0]) > 1:
                                if self.validate_pdfplumber_table(table):
                                    all_tables.append({
                                        'method': 'pdfplumber',
                                        'page': p_idx + 1,
                                        'table_index': t_idx,
                                        'data': table,
                                        'shape': (len(table), len(table[0]) if table else 0),
                                        'confidence': 'medium'
                                    })
            except:
                pass
        return all_tables

    def validate_extracted_table(self, df):
        """Validate pandas DataFrame as table"""
        if df.empty or df.shape[0] < 2 or df.shape[1] < 2:
            return False
        # Must have content
        non_empty = df.astype(str).values.flatten()
        non_empty_cells = [cell for cell in non_empty if cell.strip() and cell.strip() != 'nan']
        if len(non_empty_cells) < df.size * 0.3:
            return False
        if len(set(non_empty_cells)) < 2:
            return False
        return True

    def validate_pdfplumber_table(self, table):
        """Validate pdfplumber table"""
        if not table or len(table) < 2:
            return False
        total_cells = 0
        non_empty_cells = 0
        unique_vals = set()
        for row in table:
            for cell in row:
                total_cells += 1
                if cell and cell.strip():
                    non_empty_cells += 1
                    unique_vals.add(cell.strip())
        if non_empty_cells < total_cells * 0.3:
            return False
        if len(unique_vals) < 2:
            return False
        return True
    
    def extract_tables_optimized(self, pdf_path, max_pages=50):
        """Optimized table extraction for large files"""
        all_tables = []
        
        # For large files, use only the most reliable method and limit pages
        print(f"    📊 Optimized extraction - processing first {max_pages} pages only...")
        
        # Try pdfplumber first (most memory efficient)
        if HAS_PDFPLUMBER:
            try:
                with pdfplumber.open(pdf_path) as pdf:
                    pages_to_process = min(max_pages, len(pdf.pages))
                    for p_idx in range(pages_to_process):
                        try:
                            page = pdf.pages[p_idx]
                            tables = page.extract_tables()
                            for t_idx, table in enumerate(tables):
                                if table and len(table) > 1 and len(table[0]) > 1:
                                    if self.validate_pdfplumber_table(table):
                                        # Limit table size for memory efficiency
                                        if len(table) > 100:
                                            table = table[:100]  # Keep first 100 rows
                                        all_tables.append({
                                            'method': 'pdfplumber_optimized',
                                            'page': p_idx + 1,
                                            'table_index': t_idx,
                                            'data': table,
                                            'shape': (len(table), len(table[0]) if table else 0),
                                            'confidence': 'medium',
                                            'note': 'Limited for large file processing'
                                        })
                        except Exception as e:
                            print(f"    ⚠️  Error on page {p_idx + 1}: {str(e)}")
                            continue
                        
                        # Check memory usage periodically
                        if p_idx % 10 == 0:
                            gc.collect()
                            if self.is_memory_limit_exceeded():
                                print(f"    ⚠️  Memory limit reached at page {p_idx + 1}")
                                break
            except Exception as e:
                print(f"    ⚠️  Pdfplumber optimized error: {str(e)}")
        
        # If no tables found with pdfplumber, try tabula on first few pages only
        if not all_tables and HAS_TABULA:
            try:
                print("    📊 Trying tabula on first 10 pages...")
                pages_str = '1-10'  # Limit to first 10 pages
                tables = tabula.read_pdf(pdf_path, pages=pages_str, multiple_tables=True, silent=True)
                for i, table in enumerate(tables):
                    if not table.empty and table.shape[0] > 1 and table.shape[1] > 1:
                        if self.validate_extracted_table(table):
                            # Limit table size
                            if table.shape[0] > 100:
                                table = table.head(100)
                            all_tables.append({
                                'method': 'tabula_optimized',
                                'table_index': i,
                                'data': table.to_dict('records'),
                                'shape': table.shape,
                                'columns': table.columns.tolist(),
                                'confidence': 'high',
                                'note': 'Limited to first 10 pages and 100 rows'
                            })
            except Exception as e:
                print(f"    ⚠️  Tabula optimized error: {str(e)}")
        
        return all_tables
    
    def extract_docx(self, docx_path):
        """Extract data from DOCX files"""
        if not HAS_DOCX:
            print("❌ python-docx not available. Skipping DOCX processing.")
            return None
            
        try:
            print(f"📄 Processing DOCX: {os.path.basename(docx_path)}")
            doc = Document(docx_path)
            
            docx_data = {
                'filename': os.path.basename(docx_path),
                'file_type': 'DOCX',
                'paragraphs': [],
                'tables': [],
                'total_words': 0,
                'total_characters': 0,
                'file_size_mb': round(os.path.getsize(docx_path) / (1024*1024), 2)
            }
            
            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    docx_data['paragraphs'].append({
                        'text': para.text,
                        'style': para.style.name if para.style else 'Normal'
                    })
                    docx_data['total_words'] += len(para.text.split())
                    docx_data['total_characters'] += len(para.text)
            
            # Extract tables
            for table_idx, table in enumerate(doc.tables):
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                
                if table_data and len(table_data) > 1:
                    docx_data['tables'].append({
                        'table_index': table_idx,
                        'data': table_data,
                        'shape': (len(table_data), len(table_data[0]) if table_data else 0)
                    })
            
            print(f"  ✅ DOCX processed: {len(docx_data['paragraphs'])} paragraphs, {docx_data['total_words']} words")
            return docx_data
            
        except Exception as e:
            print(f"❌ Error processing DOCX {docx_path}: {str(e)}")
            return None
    
    def extract_pptx(self, pptx_path):
        """Extract data from PPTX files"""
        if not HAS_PPTX:
            print("❌ python-pptx not available. Skipping PPTX processing.")
            return None
            
        try:
            print(f"📄 Processing PPTX: {os.path.basename(pptx_path)}")
            prs = Presentation(pptx_path)
            
            pptx_data = {
                'filename': os.path.basename(pptx_path),
                'file_type': 'PPTX',
                'slides': [],
                'total_slides': len(prs.slides),
                'total_words': 0,
                'total_characters': 0,
                'file_size_mb': round(os.path.getsize(pptx_path) / (1024*1024), 2)
            }
            
            for slide_idx, slide in enumerate(prs.slides):
                slide_data = {
                    'slide_number': slide_idx + 1,
                    'text_content': [],
                    'tables': []
                }
                
                # Extract text from shapes
                for shape in slide.shapes:
                    try:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_data['text_content'].append(shape.text.strip())
                            pptx_data['total_words'] += len(shape.text.split())
                            pptx_data['total_characters'] += len(shape.text)
                        
                        # Extract tables
                        if hasattr(shape, "table") and shape.table:
                            table_data = []
                            for row in shape.table.rows:
                                row_data = [cell.text.strip() for cell in row.cells]
                                table_data.append(row_data)
                            if table_data and len(table_data) > 1:
                                slide_data['tables'].append(table_data)
                    except Exception as e:
                        continue
                
                pptx_data['slides'].append(slide_data)
            
            print(f"  ✅ PPTX processed: {pptx_data['total_slides']} slides, {pptx_data['total_words']} words")
            return pptx_data
            
        except Exception as e:
            print(f"❌ Error processing PPTX {pptx_path}: {str(e)}")
            return None
    
    def extract_excel(self, excel_path):
        """Extract data from Excel files with large file optimization"""
        try:
            print(f"📄 Processing Excel: {os.path.basename(excel_path)}")
            
            file_size_mb = os.path.getsize(excel_path) / (1024*1024)
            print(f"  📊 File size: {file_size_mb:.2f} MB")
            
            excel_data = {
                'filename': os.path.basename(excel_path),
                'file_type': 'Excel',
                'sheets': [],
                'file_size_mb': round(file_size_mb, 2),
                'processing_mode': 'optimized' if file_size_mb > 50 else 'standard'
            }
            
            # For large files, implement streaming and limits
            max_rows_per_sheet = 1000 if file_size_mb > 50 else None
            chunk_size = 100 if file_size_mb > 100 else None
            
            # Get sheet names without loading data
            try:
                xl_file = pd.ExcelFile(excel_path)
                sheet_names = xl_file.sheet_names
                
                for sheet_name in sheet_names:
                    try:
                        print(f"    📋 Processing sheet: {sheet_name}")
                        
                        # For large files, read in chunks or with row limits
                        if chunk_size and max_rows_per_sheet:
                            # Read with both chunking and row limits for very large files
                            print(f"    🔧 Using chunked reading with {chunk_size} rows per chunk, max {max_rows_per_sheet} rows")
                            
                            all_chunks = []
                            rows_read = 0
                            
                            try:
                                for chunk in pd.read_excel(excel_path, sheet_name=sheet_name, chunksize=chunk_size):
                                    all_chunks.append(chunk)
                                    rows_read += len(chunk)
                                    
                                    if rows_read >= max_rows_per_sheet:
                                        break
                                    
                                    # Check memory usage periodically
                                    if len(all_chunks) % 5 == 0:
                                        gc.collect()
                                        if self.is_memory_limit_exceeded():
                                            print(f"    ⚠️  Memory limit reached, stopping at {rows_read} rows")
                                            break
                                
                                # Combine chunks
                                if all_chunks:
                                    df = pd.concat(all_chunks, ignore_index=True)
                                else:
                                    df = pd.DataFrame()
                                    
                            except Exception as e:
                                print(f"    ⚠️  Chunked reading failed, trying row limit: {str(e)}")
                                # Fallback to simple row limit
                                df = pd.read_excel(excel_path, sheet_name=sheet_name, nrows=max_rows_per_sheet)
                                
                        elif max_rows_per_sheet:
                            # Read with row limit only
                            print(f"    🔧 Using row limit: {max_rows_per_sheet} rows")
                            df = pd.read_excel(excel_path, sheet_name=sheet_name, nrows=max_rows_per_sheet)
                            
                        else:
                            # Standard reading for smaller files
                            df = pd.read_excel(excel_path, sheet_name=sheet_name)
                        
                        # Process the dataframe
                        if not df.empty:
                            # For large sheets, limit the data we store
                            sample_rows = min(100, len(df)) if file_size_mb > 50 else len(df)
                            
                            sheet_data = {
                                'sheet_name': sheet_name,
                                'shape': df.shape,
                                'columns': df.columns.tolist(),
                                'data': df.head(sample_rows).to_dict('records'),
                                'data_types': df.dtypes.astype(str).to_dict()
                            }
                            
                            # Add summary statistics only for smaller files or limited rows
                            try:
                                if file_size_mb <= 50 or df.shape[0] <= 1000:
                                    sheet_data['summary'] = df.describe(include='all').to_dict()
                                else:
                                    # Limited summary for large files
                                    numeric_cols = df.select_dtypes(include=['number']).columns
                                    if len(numeric_cols) > 0:
                                        sheet_data['summary'] = df[numeric_cols].describe().to_dict()
                                    else:
                                        sheet_data['summary'] = {'note': 'Summary skipped for large non-numeric data'}
                            except Exception as e:
                                sheet_data['summary'] = {'error': f'Summary calculation failed: {str(e)}'}
                            
                            # Add processing notes for large files
                            if max_rows_per_sheet and df.shape[0] >= max_rows_per_sheet:
                                sheet_data['processing_note'] = f'Limited to first {max_rows_per_sheet} rows due to file size'
                            if file_size_mb > 50:
                                sheet_data['sample_note'] = f'Data sample limited to first {sample_rows} rows for display'
                            
                            excel_data['sheets'].append(sheet_data)
                            
                            # Clear dataframe from memory
                            del df
                            gc.collect()
                        
                        else:
                            # Empty sheet
                            sheet_data = {
                                'sheet_name': sheet_name,
                                'shape': (0, 0),
                                'columns': [],
                                'data': [],
                                'summary': {},
                                'data_types': {},
                                'processing_note': 'Sheet is empty'
                            }
                            excel_data['sheets'].append(sheet_data)
                        
                    except Exception as e:
                        print(f"    ⚠️  Error reading sheet '{sheet_name}': {str(e)}")
                        # Add error sheet info
                        error_sheet = {
                            'sheet_name': sheet_name,
                            'shape': (0, 0),
                            'columns': [],
                            'data': [],
                            'summary': {},
                            'data_types': {},
                            'processing_note': f'Error reading sheet: {str(e)}'
                        }
                        excel_data['sheets'].append(error_sheet)
                        continue
                
                # Final cleanup
                xl_file.close()
                gc.collect()
                
            except Exception as e:
                print(f"    ❌ Error opening Excel file: {str(e)}")
                return None
            
            print(f"  ✅ Excel processed: {len(excel_data['sheets'])} sheets")
            return excel_data
            
        except Exception as e:
            print(f"❌ Error processing Excel {excel_path}: {str(e)}")
            return None
    
    def process_files(self, file_paths):
        """Process all uploaded files"""
        results = {}
        
        print(f"\n🔄 Processing {len(file_paths)} file(s)...")
        print("=" * 60)
        
        for i, file_path in enumerate(file_paths, 1):
            file_ext = os.path.splitext(file_path)[1].lower()
            
            print(f"\n[{i}/{len(file_paths)}] Processing: {file_path}")
            
            try:
                if file_ext == '.pdf':
                    result = self.extract_pdf_comprehensive(file_path)
                elif file_ext == '.docx':
                    result = self.extract_docx(file_path)
                elif file_ext == '.pptx':
                    result = self.extract_pptx(file_path)
                elif file_ext in ['.xlsx', '.xls']:
                    result = self.extract_excel(file_path)
                else:
                    print(f"❌ Unsupported file format: {file_ext}")
                    continue
                
                if result:
                    results[file_path] = result
                    print(f"✅ Successfully processed: {file_path}")
                else:
                    print(f"❌ Failed to process: {file_path}")
                    
            except Exception as e:
                print(f"❌ Error processing {file_path}: {str(e)}")
                continue
        
        return results
    
    def save_results(self, results, output_format='json'):
        """Save extraction results with better formatting"""
        if not results:
            print("❌ No results to save.")
            return
        
        try:
            # Save as JSON
            with open('extraction_results.json', 'w', encoding='utf-8') as f:
                json.dump(results, f, indent=2, ensure_ascii=False, default=str)
            print("✅ Results saved as: extraction_results.json")
            
            # Create summary report
            self.create_summary_report(results)
            
            # Create CSV summaries for tables
            self.save_tables_as_csv(results)
            
            # Check if there are extracted images
            if os.path.exists('extracted_images') and os.listdir('extracted_images'):
                print("🖼️  Images extracted to 'extracted_images' folder")
                
        except Exception as e:
            print(f"❌ Error saving results: {str(e)}")
    
    def save_tables_as_csv(self, results):
        """Save extracted tables as CSV files"""
        table_count = 0
        
        for file_path, data in results.items():
            filename_base = os.path.splitext(data['filename'])[0]
            
            # Save tables from different extraction methods
            if 'extracted_tables' in data:
                for i, table in enumerate(data['extracted_tables']):
                    try:
                        if 'data' in table and table['data']:
                            csv_filename = f"extracted_tables/{filename_base}_table_{i+1}_{table['method']}.csv"
                            
                            if table['method'] in ['tabula', 'camelot_lattice']:
                                # These are already in records format
                                df = pd.DataFrame(table['data'])
                            else:
                                # Convert list of lists to DataFrame
                                df = pd.DataFrame(table['data'])
                            
                            df.to_csv(csv_filename, index=False)
                            table_count += 1
                            
                    except Exception as e:
                        print(f"    ⚠️  Error saving table {i+1}: {str(e)}")
                        continue
        
        if table_count > 0:
            print(f"📊 Saved {table_count} tables as CSV files in 'extracted_tables' folder")
    
    def create_summary_report(self, results):
        """Create a comprehensive summary report"""
        try:
            with open('extraction_summary.txt', 'w', encoding='utf-8') as f:
                f.write("📋 DOCUMENT EXTRACTION SUMMARY REPORT\n")
                f.write("=" * 60 + "\n")
                f.write(f"📅 Generated: {pd.Timestamp.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
                f.write(f"📁 Total Files Processed: {len(results)}\n\n")
                
                # Overall statistics
                total_words = sum(data.get('total_words', 0) for data in results.values())
                total_chars = sum(data.get('total_characters', 0) for data in results.values())
                total_images = sum(data.get('total_images', 0) for data in results.values())
                total_tables = sum(len(data.get('extracted_tables', [])) for data in results.values())
                total_size_mb = sum(data.get('file_size_mb', 0) for data in results.values())
                
                f.write("📊 OVERALL STATISTICS:\n")
                f.write("-" * 30 + "\n")
                f.write(f"📝 Total Words: {total_words:,}\n")
                f.write(f"🔤 Total Characters: {total_chars:,}\n")
                f.write(f"🖼️  Total Images: {total_images:,}\n")
                f.write(f"📊 Total Tables: {total_tables:,}\n")
                f.write(f"💾 Total File Size: {total_size_mb:.2f} MB\n\n")
                
                # File type distribution
                file_types = Counter(data['file_type'] for data in results.values())
                f.write("📈 FILE TYPE DISTRIBUTION:\n")
                f.write("-" * 30 + "\n")
                for file_type, count in file_types.items():
                    f.write(f"{file_type}: {count} files\n")
                f.write("\n")
                
                # Detailed file information
                f.write("📄 DETAILED FILE INFORMATION:\n")
                f.write("=" * 60 + "\n")
                
                for file_path, data in results.items():
                    f.write(f"\n📁 File: {data['filename']}\n")
                    f.write(f"   Type: {data['file_type']}\n")
                    f.write(f"   Size: {data.get('file_size_mb', 0):.2f} MB\n")
                    
                    if data['file_type'] == 'PDF':
                        f.write(f"   Pages: {data.get('page_count', 0)}\n")
                        f.write(f"   Words: {data.get('total_words', 0):,}\n")
                        f.write(f"   Characters: {data.get('total_characters', 0):,}\n")
                        f.write(f"   Images: {data.get('total_images', 0)}\n")
                        f.write(f"   Tables: {len(data.get('extracted_tables', []))}\n")
                        
                        # Font information
                        fonts = data.get('fonts_used', [])
                        if fonts:
                            f.write(f"   Fonts Used: {', '.join(fonts[:5])}")
                            if len(fonts) > 5:
                                f.write(f" ... (+{len(fonts)-5} more)")
                            f.write("\n")
                        
                        # Metadata
                        metadata = data.get('metadata', {})
                        if metadata:
                            f.write("   Metadata:\n")
                            for key, value in metadata.items():
                                f.write(f"     {key}: {value}\n")
                    
                    elif data['file_type'] == 'DOCX':
                        f.write(f"   Paragraphs: {len(data.get('paragraphs', []))}\n")
                        f.write(f"   Words: {data.get('total_words', 0):,}\n")
                        f.write(f"   Characters: {data.get('total_characters', 0):,}\n")
                        f.write(f"   Tables: {len(data.get('tables', []))}\n")
                    
                    elif data['file_type'] == 'PPTX':
                        f.write(f"   Slides: {data.get('total_slides', 0)}\n")
                        f.write(f"   Words: {data.get('total_words', 0):,}\n")
                        f.write(f"   Characters: {data.get('total_characters', 0):,}\n")
                        
                        # Count tables across all slides
                        total_slide_tables = sum(len(slide.get('tables', [])) for slide in data.get('slides', []))
                        f.write(f"   Tables: {total_slide_tables}\n")
                    
                    elif data['file_type'] == 'Excel':
                        f.write(f"   Sheets: {len(data.get('sheets', []))}\n")
                        
                        # Sheet details
                        for sheet in data.get('sheets', []):
                            f.write(f"     - {sheet['sheet_name']}: {sheet['shape'][0]} rows x {sheet['shape'][1]} cols\n")
                
                # Table extraction summary
                if total_tables > 0:
                    f.write(f"\n📊 TABLE EXTRACTION SUMMARY:\n")
                    f.write("-" * 40 + "\n")
                    
                    method_counts = Counter()
                    confidence_counts = Counter()
                    
                    for data in results.values():
                        for table in data.get('extracted_tables', []):
                            method_counts[table.get('method', 'unknown')] += 1
                            confidence_counts[table.get('confidence', 'unknown')] += 1
                    
                    f.write("Extraction Methods:\n")
                    for method, count in method_counts.items():
                        f.write(f"  {method}: {count} tables\n")
                    
                    f.write("\nConfidence Levels:\n")
                    for confidence, count in confidence_counts.items():
                        f.write(f"  {confidence}: {count} tables\n")
                
                # Image extraction summary
                if total_images > 0:
                    f.write(f"\n🖼️  IMAGE EXTRACTION SUMMARY:\n")
                    f.write("-" * 40 + "\n")
                    
                    for data in results.values():
                        if data.get('total_images', 0) > 0:
                            f.write(f"From {data['filename']}: {data['total_images']} images\n")
                            
                            # Image details from pages
                            for page in data.get('pages', []):
                                if page.get('images'):
                                    f.write(f"  Page {page['page_number']}: {len(page['images'])} images\n")
                
                # Processing warnings and recommendations
                f.write(f"\n⚠️  PROCESSING NOTES:\n")
                f.write("-" * 30 + "\n")
                
                # Check for potential issues
                large_files = [data
Thought for 1s
Let me continue with the complete code files:

pdfextract.py (continued from where it was cut off)
                # Check for potential issues
                large_files = [data for data in results.values() if data.get('file_size_mb', 0) > 50]
                if large_files:
                    f.write(f"📁 Large files detected ({len(large_files)} files > 50MB)\n")
                
                # Check for files with many images
                image_heavy_files = [data for data in results.values() if data.get('total_images', 0) > 20]
                if image_heavy_files:
                    f.write(f"🖼️  Image-heavy files detected ({len(image_heavy_files)} files > 20 images)\n")
                
                # Check for missing libraries
                missing_libs = []
                if not HAS_PDFPLUMBER:
                    missing_libs.append("pdfplumber")
                if not HAS_TABULA:
                    missing_libs.append("tabula-py")
                if not HAS_CAMELOT:
                    missing_libs.append("camelot-py")
                if not HAS_CV2:
                    missing_libs.append("opencv-python")
                
                if missing_libs:
                    f.write(f"📦 Missing optional libraries: {', '.join(missing_libs)}\n")
                    f.write("   Install these for enhanced extraction capabilities\n")
                
                f.write(f"\n✅ Report generated successfully!\n")
                f.write("   Files created:\n")
                f.write("   - extraction_results.json (full data)\n")
                f.write("   - extraction_summary.txt (this report)\n")
                f.write("   - extracted_images/ (folder with images)\n")
                f.write("   - extracted_tables/ (folder with CSV tables)\n")
            
            print("✅ Summary report saved as: extraction_summary.txt")
            
        except Exception as e:
            print(f"❌ Error creating summary report: {str(e)}")
    
    def generate_extraction_stats(self, results):
        """Generate detailed extraction statistics"""
        stats = {
            'total_files': len(results),
            'file_types': Counter(data['file_type'] for data in results.values()),
            'total_words': sum(data.get('total_words', 0) for data in results.values()),
            'total_characters': sum(data.get('total_characters', 0) for data in results.values()),
            'total_images': sum(data.get('total_images', 0) for data in results.values()),
            'total_tables': sum(len(data.get('extracted_tables', [])) for data in results.values()),
            'total_size_mb': sum(data.get('file_size_mb', 0) for data in results.values()),
            'processing_time': getattr(self, 'processing_time', 0)
        }
        
        # Add file-specific stats
        stats['pdf_stats'] = {
            'count': sum(1 for data in results.values() if data['file_type'] == 'PDF'),
            'total_pages': sum(data.get('page_count', 0) for data in results.values() if data['file_type'] == 'PDF'),
            'fonts_used': list(set(font for data in results.values() if data['file_type'] == 'PDF' for font in data.get('fonts_used', [])))
        }
        
        return stats
    
    def cleanup_temp_files(self):
        """Clean up temporary files and directories"""
        try:
            # Remove empty directories
            for folder in ['extracted_images', 'extracted_tables']:
                if os.path.exists(folder) and not os.listdir(folder):
                    os.rmdir(folder)
                    print(f"🧹 Removed empty folder: {folder}")
        except Exception as e:
            print(f"⚠️  Error during cleanup: {str(e)}")

def main():
    """Main function to run the document extractor"""
    print("🚀 Comprehensive Document Extractor")
    print("=" * 60)
    print("Supported formats: PDF, DOCX, PPTX, XLSX, XLS")
    print()
    
    # Initialize the extractor
    extractor = ComprehensiveDocumentExtractor()
    
    # Get file paths from user
    print("📁 Please provide file paths to process:")
    print("   - Enter file paths separated by commas")
    print("   - Or drag and drop files into the terminal")
    print("   - Press Enter when done")
    print()
    
    file_paths = []
    while True:
        user_input = input("File path (or 'done' to finish): ").strip()
        
        if user_input.lower() in ['done', 'finish', 'exit', '']:
            break
        
        # Handle multiple files separated by commas
        if ',' in user_input:
            paths = [path.strip().strip('"\'') for path in user_input.split(',')]
            file_paths.extend(paths)
        else:
            file_paths.append(user_input.strip('"\''))
    
    if not file_paths:
        print("❌ No files provided. Exiting.")
        return
    
    # Validate file paths
    valid_paths = []
    for path in file_paths:
        if os.path.exists(path):
            ext = os.path.splitext(path)[1].lower()
            if ext in extractor.supported_formats:
                valid_paths.append(path)
            else:
                print(f"⚠️  Skipping unsupported file: {path}")
        else:
            print(f"⚠️  File not found: {path}")
    
    if not valid_paths:
        print("❌ No valid files to process. Exiting.")
        return
    
    # Record start time
    start_time = pd.Timestamp.now()
    
    # Process files
    try:
        results = extractor.process_files(valid_paths)
        
        if results:
            # Record processing time
            extractor.processing_time = (pd.Timestamp.now() - start_time).total_seconds()
            
            # Save results
            extractor.save_results(results)
            
            # Generate and display final statistics
            stats = extractor.generate_extraction_stats(results)
            
            print("\n" + "=" * 60)
            print("📊 FINAL EXTRACTION STATISTICS")
            print("=" * 60)
            print(f"✅ Successfully processed: {stats['total_files']} files")
            print(f"📝 Total words extracted: {stats['total_words']:,}")
            print(f"🔤 Total characters: {stats['total_characters']:,}")
            print(f"🖼️  Total images: {stats['total_images']:,}")
            print(f"📊 Total tables: {stats['total_tables']:,}")
            print(f"💾 Total file size: {stats['total_size_mb']:.2f} MB")
            print(f"⏱️  Processing time: {stats['processing_time']:.2f} seconds")
            
            if stats['pdf_stats']['count'] > 0:
                print(f"📄 PDF files: {stats['pdf_stats']['count']} ({stats['pdf_stats']['total_pages']} pages)")
            
            print("\n📁 Output files created:")
            print("   - extraction_results.json")
            print("   - extraction_summary.txt")
            if stats['total_images'] > 0:
                print("   - extracted_images/ (folder)")
            if stats['total_tables'] > 0:
                print("   - extracted_tables/ (folder)")
            
            # Cleanup
            extractor.cleanup_temp_files()
            
            print("\n🎉 Document extraction completed successfully!")
            
        else:
            print("❌ No files were successfully processed.")
            
    except KeyboardInterrupt:
        print("\n⏹️  Process interrupted by user.")
        extractor.cleanup_temp_files()
    except Exception as e:
        print(f"\n❌ Unexpected error: {str(e)}")
        extractor.cleanup_temp_files()

if __name__ == "__main__":
    main()
