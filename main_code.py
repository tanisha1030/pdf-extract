import fitz  # PyMuPDF
import pandas as pd
from docx import Document
from pptx import Presentation
from PIL import Image
from io import BytesIO
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_from_pdf(file):
    doc = fitz.open(stream=file.read(), filetype="pdf")
    pages_info = []
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        text = page.get_text()
        images = []
        for img_index, img in enumerate(page.get_images(full=True)):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            images.append(Image.open(BytesIO(image_bytes)))
        tables = page.find_tables()
        tables_data = []
        for table in tables:
            try:
                df = pd.DataFrame(table.extract())
                tables_data.append(df)
            except:
                pass
        pages_info.append({"text": text, "images": images, "tables": tables_data})
    return pages_info

def extract_from_docx(file):
    doc = Document(file)
    text = "\n".join([para.text for para in doc.paragraphs])
    return {"text": text}

def extract_table_from_pptx_shape(table_shape):
    """Extract table data from a PowerPoint table shape"""
    try:
        table = table_shape.table
        data = []
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                row_data.append(cell_text)
            data.append(row_data)
        
        if data:
            # Create DataFrame with first row as headers if it looks like headers
            df = pd.DataFrame(data[1:], columns=data[0]) if len(data) > 1 else pd.DataFrame(data)
            return df
        return None
    except Exception as e:
        print(f"Error extracting table: {e}")
        return None

def extract_from_pptx(file):
    prs = Presentation(file)
    slides_data = []
    
    for slide_num, slide in enumerate(prs.slides):
        text = []
        tables = []
        images = []
        
        for shape in slide.shapes:
            try:
                # Extract text from text boxes and other text-containing shapes
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text.strip())
                
                # Extract tables using the correct shape type check
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    print(f"Found table in slide {slide_num + 1}")
                    table_df = extract_table_from_pptx_shape(shape)
                    if table_df is not None and not table_df.empty:
                        tables.append(table_df)
                
                # Extract images using the correct shape type check
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    print(f"Found image in slide {slide_num + 1}")
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        pil_image = Image.open(BytesIO(image_bytes))
                        images.append(pil_image)
                    except Exception as img_error:
                        print(f"Error extracting image: {img_error}")
                        continue
                
                # Also check for grouped shapes that might contain tables/images
                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    print(f"Found group in slide {slide_num + 1}")
                    for grouped_shape in shape.shapes:
                        if hasattr(grouped_shape, "text") and grouped_shape.text.strip():
                            text.append(grouped_shape.text.strip())
                        
                        if grouped_shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                            table_df = extract_table_from_pptx_shape(grouped_shape)
                            if table_df is not None and not table_df.empty:
                                tables.append(table_df)
                        
                        elif grouped_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            try:
                                image = grouped_shape.image
                                image_bytes = image.blob
                                pil_image = Image.open(BytesIO(image_bytes))
                                images.append(pil_image)
                            except Exception as img_error:
                                print(f"Error extracting grouped image: {img_error}")
                                continue
            
            except Exception as shape_error:
                print(f"Error processing shape in slide {slide_num + 1}: {shape_error}")
                continue
        
        print(f"Slide {slide_num + 1}: Found {len(tables)} tables and {len(images)} images")
        
        slides_data.append({
            "text": "\n".join(text),
            "tables": tables,
            "images": images
        })
    
    return slides_data

def extract_from_excel(file):
    xls = pd.ExcelFile(file)
    sheets_data = {}
    for sheet_name in xls.sheet_names:
        df = pd.read_excel(xls, sheet_name)
        sheets_data[sheet_name] = df
    return sheets_data
