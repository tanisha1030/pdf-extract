import fitz  # PyMuPDF
import io
import os
import json
import pandas as pd
from PIL import Image
import pytesseract
from collections import Counter
import zipfile
import tempfile
from google.colab import files
import pdfplumber
import tabula
import camelot
import cv2
import numpy as np
from docx import Document
from pptx import Presentation
import openpyxl
import warnings
warnings.filterwarnings('ignore')

class ComprehensiveDocumentExtractor:
    def __init__(self):
        self.extracted_data = {}
        self.supported_formats = ['.pdf', '.docx', '.pptx', '.xlsx', '.xls']
        self.setup_directories()
    
    def setup_directories(self):
        """Create necessary directories"""
        os.makedirs('extracted_images', exist_ok=True)
        os.makedirs('extracted_tables', exist_ok=True)
    
    def upload_files(self):
        """Allow user to upload multiple files with better error handling"""
        print("📤 Please select your files to upload...")
        print(f"📋 Supported formats: {', '.join(self.supported_formats)}")
        
        try:
            uploaded = files.upload()
            
            if not uploaded:
                print("❌ No files uploaded.")
                return []
            
            file_paths = []
            for filename, content in uploaded.items():
                file_ext = os.path.splitext(filename)[1].lower()
                
                if file_ext not in self.supported_formats:
                    print(f"⚠️  Skipping unsupported file: {filename} ({file_ext})")
                    continue
                
                # Save uploaded file
                try:
                    with open(filename, 'wb') as f:
                        f.write(content)
                    file_paths.append(filename)
                    print(f"✅ Uploaded: {filename} ({len(content):,} bytes)")
                except Exception as e:
                    print(f"❌ Failed to save {filename}: {str(e)}")
            
            return file_paths
            
        except Exception as e:
            print(f"❌ Upload failed: {str(e)}")
            return []
    
    def extract_pdf_comprehensive(self, pdf_path):
        """Extract comprehensive data from PDF with error handling"""
        try:
            print(f"📄 Processing PDF: {os.path.basename(pdf_path)}")
            doc = fitz.open(pdf_path)
            
            pdf_data = {
                'filename': os.path.basename(pdf_path),
                'file_type': 'PDF',
                'metadata': {},
                'pages': [],
                'total_images': 0,
                'total_tables': 0,
                'fonts_used': [],
                'total_words': 0,
                'total_characters': 0,
                'page_count': len(doc),
                'file_size_mb': round(os.path.getsize(pdf_path) / (1024*1024), 2)
            }
            
            # Extract metadata safely
            try:
                metadata = doc.metadata
                pdf_data['metadata'] = {k: v for k, v in metadata.items() if v}
            except:
                pdf_data['metadata'] = {}
            
            # Process each page with progress indicator
            for page_num in range(len(doc)):
                if page_num % 10 == 0:
                    print(f"  📃 Processing page {page_num + 1}/{len(doc)}")
                
                try:
                    page = doc[page_num]
                    page_data = self.extract_page_data(page, page_num + 1, pdf_path)
                    pdf_data['pages'].append(page_data)
                    
                    # Accumulate totals
                    pdf_data['total_images'] += len(page_data['images'])
                    pdf_data['total_tables'] += len(page_data['tables'])
                    pdf_data['total_words'] += page_data['word_count']
                    pdf_data['total_characters'] += page_data['char_count']
                    
                except Exception as e:
                    print(f"  ⚠️  Error on page {page_num + 1}: {str(e)}")
                    continue
            
            # Extract unique fonts
            all_fonts = []
            for page in pdf_data['pages']:
                all_fonts.extend(page.get('fonts', []))
            pdf_data['fonts_used'] = list(set(all_fonts))
            
            # Extract tables using multiple methods
            print("  📊 Extracting tables...")
            pdf_data['extracted_tables'] = self.extract_tables_multiple_methods(pdf_path)
            
            doc.close()
            print(f"  ✅ PDF processed: {pdf_data['page_count']} pages, {pdf_data['total_words']} words")
            return pdf_data
            
        except Exception as e:
            print(f"❌ Error processing PDF {pdf_path}: {str(e)}")
            return None
    
    def extract_page_data(self, page, page_num, pdf_path):
        """Extract comprehensive data from a single page"""
        page_data = {
            'page_number': page_num,
            'text': '',
            'formatted_text': [],
            'images': [],
            'tables': [],
            'fonts': [],
            'word_count': 0,
            'char_count': 0,
            'page_size': {
                'width': round(page.rect.width, 2),
                'height': round(page.rect.height, 2)
            },
            'rotation': page.rotation
        }
        
        try:
            # Extract text with formatting
            text_dict = page.get_text("dict")
            page_text = ""
            fonts_on_page = []
            
            for block in text_dict.get("blocks", []):
                if "lines" in block:
                    for line in block["lines"]:
                        for span in line["spans"]:
                            text = span.get("text", "")
                            page_text += text
                            
                            # Store formatted text with details
                            page_data['formatted_text'].append({
                                'text': text,
                                'font': span.get('font', 'Unknown'),
                                'size': round(span.get('size', 0), 2),
                                'flags': span.get('flags', 0),
                                'color': span.get('color', 0),
                                'bbox': span.get('bbox', [0, 0, 0, 0])
                            })
                            
                            fonts_on_page.append(span.get('font', 'Unknown'))
            
            page_data['text'] = page_text
            page_data['fonts'] = list(set(fonts_on_page))
            page_data['word_count'] = len(page_text.split())
            page_data['char_count'] = len(page_text)
            
            # Extract images
            page_data['images'] = self.extract_images_from_page(page, page_num, pdf_path)
            
            # Extract tables (basic detection)
            page_data['tables'] = self.extract_tables_from_page(page)
            
        except Exception as e:
            print(f"    ⚠️  Error extracting from page {page_num}: {str(e)}")
        
        return page_data
    
    def extract_images_from_page(self, page, page_num, pdf_path):
        """Extract images from a page with better error handling"""
        images = []
        
        try:
            image_list = page.get_images()
            
            for img_index, img in enumerate(image_list):
                try:
                    xref = img[0]
                    pix = fitz.Pixmap(page.parent, xref)
                    
                    if pix.n - pix.alpha < 4:  # GRAY or RGB
                        img_data = pix.tobytes("png")
                        
                        # Create safe filename
                        base_name = os.path.splitext(os.path.basename(pdf_path))[0]
                        img_name = f"extracted_images/{base_name}_page_{page_num}_img_{img_index + 1}.png"
                        
                        # Save image
                        with open(img_name, "wb") as img_file:
                            img_file.write(img_data)
                        
                        # Get image info
                        img_info = {
                            'filename': img_name,
                            'width': pix.width,
                            'height': pix.height,
                            'colorspace': pix.colorspace.name if pix.colorspace else 'Unknown',
                            'size_bytes': len(img_data)
                        }
                        
                        images.append(img_info)
                    
                    pix = None
                    
                except Exception as e:
                    print(f"      ⚠️  Error extracting image {img_index} from page {page_num}: {str(e)}")
                    continue
                    
        except Exception as e:
            print(f"    ⚠️  Error accessing images on page {page_num}: {str(e)}")
        
        return images
    
    def extract_tables_from_page(self, page):
        """Basic table detection from page"""
        tables = []
        try:
            text_dict = page.get_text("dict")
            
            for block in text_dict.get("blocks", []):
                if "lines" in block:
                    lines = block["lines"]
                    if len(lines) > 2:  # Potential table
                        table_data = []
                        for line in lines:
                            row_text = ""
                            for span in line["spans"]:
                                row_text += span.get("text", "") + " "
                            table_data.append(row_text.strip())
                        
                        if len(table_data) > 2:
                            tables.append({
                                'bbox': block.get('bbox', [0, 0, 0, 0]),
                                'data': table_data,
                                'method': 'basic_detection'
                            })
            
        except Exception as e:
            print(f"      ⚠️  Error extracting tables from page: {str(e)}")
        
        return tables
    
    def extract_tables_multiple_methods(self, pdf_path):
        """Extract tables using multiple libraries for better accuracy"""
        all_tables = []
        
        # Method 1: Using tabula-py
        try:
            print("    📊 Trying tabula...")
            tabula_tables = tabula.read_pdf(pdf_path, pages='all', multiple_tables=True, silent=True)
            for i, table in enumerate(tabula_tables):
                if not table.empty:
                    all_tables.append({
                        'method': 'tabula',
                        'table_index': i,
                        'data': table.to_dict('records'),
                        'shape': table.shape,
                        'columns': table.columns.tolist()
                    })
            print(f"      ✅ Tabula found {len(tabula_tables)} tables")
        except Exception as e:
            print(f"      ⚠️  Tabula extraction failed: {str(e)}")
        
        # Method 2: Using camelot (for lattice tables)
        try:
            print("    📊 Trying camelot...")
            camelot_tables = camelot.read_pdf(pdf_path, pages='all', flavor='lattice')
            for i, table in enumerate(camelot_tables):
                if not table.df.empty:
                    all_tables.append({
                        'method': 'camelot_lattice',
                        'table_index': i,
                        'data': table.df.to_dict('records'),
                        'shape': table.df.shape,
                        'accuracy': getattr(table, 'accuracy', 'N/A')
                    })
            print(f"      ✅ Camelot found {len(camelot_tables)} tables")
        except Exception as e:
            print(f"      ⚠️  Camelot extraction failed: {str(e)}")
        
        # Method 3: Using pdfplumber
        try:
            print("    📊 Trying pdfplumber...")
            table_count = 0
            with pdfplumber.open(pdf_path) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    tables = page.extract_tables()
                    for i, table in enumerate(tables):
                        if table and len(table) > 1:
                            all_tables.append({
                                'method': 'pdfplumber',
                                'page': page_num + 1,
                                'table_index': i,
                                'data': table,
                                'shape': (len(table), len(table[0]) if table else 0)
                            })
                            table_count += 1
            print(f"      ✅ PDFplumber found {table_count} tables")
        except Exception as e:
            print(f"      ⚠️  PDFplumber extraction failed: {str(e)}")
        
        return all_tables
    
    def extract_docx(self, docx_path):
        """Extract data from DOCX files"""
        try:
            print(f"📄 Processing DOCX: {os.path.basename(docx_path)}")
            doc = Document(docx_path)
            
            docx_data = {
                'filename': os.path.basename(docx_path),
                'file_type': 'DOCX',
                'paragraphs': [],
                'tables': [],
                'total_words': 0,
                'total_characters': 0,
                'file_size_mb': round(os.path.getsize(docx_path) / (1024*1024), 2)
            }
            
            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    docx_data['paragraphs'].append({
                        'text': para.text,
                        'style': para.style.name if para.style else 'Normal'
                    })
                    docx_data['total_words'] += len(para.text.split())
                    docx_data['total_characters'] += len(para.text)
            
            # Extract tables
            for table_idx, table in enumerate(doc.tables):
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                
                if table_data:
                    docx_data['tables'].append({
                        'table_index': table_idx,
                        'data': table_data,
                        'shape': (len(table_data), len(table_data[0]) if table_data else 0)
                    })
            
            print(f"  ✅ DOCX processed: {len(docx_data['paragraphs'])} paragraphs, {docx_data['total_words']} words")
            return docx_data
            
        except Exception as e:
            print(f"❌ Error processing DOCX {docx_path}: {str(e)}")
            return None
    
    def extract_pptx(self, pptx_path):
        """Extract data from PPTX files"""
        try:
            print(f"📄 Processing PPTX: {os.path.basename(pptx_path)}")
            prs = Presentation(pptx_path)
            
            pptx_data = {
                'filename': os.path.basename(pptx_path),
                'file_type': 'PPTX',
                'slides': [],
                'total_slides': len(prs.slides),
                'total_words': 0,
                'total_characters': 0,
                'file_size_mb': round(os.path.getsize(pptx_path) / (1024*1024), 2)
            }
            
            for slide_idx, slide in enumerate(prs.slides):
                slide_data = {
                    'slide_number': slide_idx + 1,
                    'text_content': [],
                    'tables': []
                }
                
                # Extract text from shapes
                for shape in slide.shapes:
                    try:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_data['text_content'].append(shape.text.strip())
                            pptx_data['total_words'] += len(shape.text.split())
                            pptx_data['total_characters'] += len(shape.text)
                        
                        # Extract tables
                        if hasattr(shape, "table") and shape.table:
                            table_data = []
                            for row in shape.table.rows:
                                row_data = [cell.text.strip() for cell in row.cells]
                                table_data.append(row_data)
                            if table_data:
                                slide_data['tables'].append(table_data)
                    except Exception as e:
                        continue
                
                pptx_data['slides'].append(slide_data)
            
            print(f"  ✅ PPTX processed: {pptx_data['total_slides']} slides, {pptx_data['total_words']} words")
            return pptx_data
            
        except Exception as e:
            print(f"❌ Error processing PPTX {pptx_path}: {str(e)}")
            return None
    
    def extract_excel(self, excel_path):
        """Extract data from Excel files"""
        try:
            print(f"📄 Processing Excel: {os.path.basename(excel_path)}")
            
            excel_data = {
                'filename': os.path.basename(excel_path),
                'file_type': 'Excel',
                'sheets': [],
                'file_size_mb': round(os.path.getsize(excel_path) / (1024*1024), 2)
            }
            
            # Read all sheets
            xl_file = pd.ExcelFile(excel_path)
            
            for sheet_name in xl_file.sheet_names:
                try:
                    df = pd.read_excel(excel_path, sheet_name=sheet_name)
                    
                    sheet_data = {
                        'sheet_name': sheet_name,
                        'shape': df.shape,
                        'columns': df.columns.tolist(),
                        'data': df.head(100).to_dict('records') if not df.empty else [],  # Limit to first 100 rows
                        'summary': df.describe(include='all').to_dict() if not df.empty else {},
                        'data_types': df.dtypes.astype(str).to_dict()
                    }
                    
                    excel_data['sheets'].append(sheet_data)
                    
                except Exception as e:
                    print(f"    ⚠️  Error reading sheet '{sheet_name}': {str(e)}")
                    continue
            
            print(f"  ✅ Excel processed: {len(excel_data['sheets'])} sheets")
            return excel_data
            
        except Exception as e:
            print(f"❌ Error processing Excel {excel_path}: {str(e)}")
            return None
    
    def process_files(self, file_paths):
        """Process all uploaded files"""
        results = {}
        
        print(f"\n🔄 Processing {len(file_paths)} file(s)...")
        print("=" * 60)
        
        for i, file_path in enumerate(file_paths, 1):
            file_ext = os.path.splitext(file_path)[1].lower()
            
            print(f"\n[{i}/{len(file_paths)}] Processing: {file_path}")
            
            try:
                if file_ext == '.pdf':
                    result = self.extract_pdf_comprehensive(file_path)
                elif file_ext == '.docx':
                    result = self.extract_docx(file_path)
                elif file_ext == '.pptx':
                    result = self.extract_pptx(file_path)
                elif file_ext in ['.xlsx', '.xls']:
                    result = self.extract_excel(file_path)
                else:
                    print(f"❌ Unsupported file format: {file_ext}")
                    continue
                
                if result:
                    results[file_path] = result
                    print(f"✅ Successfully processed: {file_path}")
                else:
                    print(f"❌ Failed to process: {file_path}")
                    
            except Exception as e:
                print(f"❌ Error processing {file_path}: {str(e)}")
                continue
        
        return results
    
    def save_results(self, results, output_format='json'):
        """Save extraction results with better formatting"""
        if not results:
            print("❌ No results to save.")
            return
        
        try:
            # Save as JSON
            with open('extraction_results.json', 'w', encoding='utf-8') as f:
                json.dump(results, f, indent=2, ensure_ascii=False, default=str)
            print("✅ Results saved as: extraction_results.json")
            
            # Create summary report
            self.create_summary_report(results)
            
            # Create CSV summaries for tables
            self.save_tables_as_csv(results)
            
            # Download results
            print("\n📥 Downloading files...")
            files.download('extraction_results.json')
            files.download('extraction_summary.txt')
            
            # Check if there are extracted images
            if os.path.exists('extracted_images') and os.listdir('extracted_images'):
                print("🖼️  Images extracted to 'extracted_images' folder")
                
        except Exception as e:
            print(f"❌ Error saving results: {str(e)}")
    
    def save_tables_as_csv(self, results):
        """Save extracted tables as CSV files"""
        table_count = 0
        
        for file_path, data in results.items():
            filename_base = os.path.splitext(data['filename'])[0]
            
            # Save tables from different extraction methods
            if 'extracted_tables' in data:
                for i, table in enumerate(data['extracted_tables']):
                    try:
                        if 'data' in table and table['data']:
                            csv_filename = f"extracted_tables/{filename_base}_table_{i+1}_{table['method']}.csv"
                            
                            if table['method'] in ['tabula', 'camelot_lattice']:
                                # These are already in records format
                                df = pd.DataFrame(table['data'])
                            else:
                                # Convert list of lists to DataFrame
                                df = pd.DataFrame(table['data'])
                            
                            df.to_csv(csv_filename, index=False)
                            table_count += 1
                            
                    except Exception as e:
                        print(f"    ⚠️  Error saving table {i+1}: {str(e)}")
                        continue
        
        if table_count > 0:
            print(f"📊 Saved {table_count} tables as CSV files in 'extracted_tables' folder")
    
    def create_summary_report(self, results):
        """Create a comprehensive summary report"""
        try:
            with open('extraction_summary.txt', 'w', encoding='utf-8') as f:
                f.write("📋 DOCUMENT EXTRACTION SUMMARY REPORT\n")
                f.write("=" * 60 + "\n")
                f.write(f"📅 Generated: {pd.Timestamp.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
                f.write(f"📁 Total Files Processed: {len(results)}\n\n")
                
                # Overall statistics
                total_words = sum(data.get('total_words', 0) for data in results.values())
                total_chars = sum(data.get('total_characters', 0) for data in results.values())
                total_images = sum(data.get('total_images', 0) for data in results.values())
                
                f.write("📊 OVERALL STATISTICS:\n")
                f.write("-" * 30 + "\n")
                f.write(f"📝 Total Words: {total_words:,}\n")
                f.write(f"🔤 Total Characters: {total_chars:,}\n")
                f.write(f"🖼️ Total Images: {total_images}\n\n")
                
                # File-by-file details
                for file_path, data in results.items():
                    f.write("=" * 60 + "\n")
                    f.write(f"📄 FILE: {data.get('filename', file_path)}\n")
                    f.write("=" * 60 + "\n")
                    
                    file_type = data.get('file_type', 'Unknown')
                    f.write(f"📊 Type: {file_type}\n")
                    f.write(f"💾 Size: {data.get('file_size_mb', 0)} MB\n")
                    
                    if file_type == 'PDF':
                        f.write(f"📃 Pages: {data.get('page_count', 0)}\n")
                        f.write(f"📝 Words: {data.get('total_words', 0):,}\n")
                        f.write(f"🔤 Characters: {data.get('total_characters', 0):,}\n")
                        f.write(f"🖼️ Images: {data.get('total_images', 0)}\n")
                        f.write(f"📋 Tables: {len(data.get('extracted_tables', []))}\n")
                        
                        fonts = data.get('fonts_used', [])
                        if fonts:
                            f.write(f"🔤 Fonts: {', '.join(fonts[:5])}")
                            if len(fonts) > 5:
                                f.write(f" ... and {len(fonts)-5} more")
                            f.write("\n")
                        
                        # Metadata
                        metadata = data.get('metadata', {})
                        if metadata:
                            f.write(f"\n📋 METADATA:\n")
                            for key, value in metadata.items():
                                if value:
                                    f.write(f"  {key}: {value}\n")
                    
                    elif file_type == 'PPTX':
                        f.write(f"📃 Slides: {data.get('total_slides', 0)}\n")
                        f.write(f"📝 Words: {data.get('total_words', 0):,}\n")
                        f.write(f"🔤 Characters: {data.get('total_characters', 0):,}\n")
                    
                    elif file_type == 'Excel':
                        sheets = data.get('sheets', [])
                        f.write(f"📃 Sheets: {len(sheets)}\n")
                        for sheet in sheets:
                            f.write(f"  - {sheet['sheet_name']}: {sheet['shape'][0]:,} rows × {sheet['shape'][1]} columns\n")
                    
                    elif file_type == 'DOCX':
                        f.write(f"📃 Paragraphs: {len(data.get('paragraphs', []))}\n")
                        f.write(f"📝 Words: {data.get('total_words', 0):,}\n")
                        f.write(f"🔤 Characters: {data.get('total_characters', 0):,}\n")
                        f.write(f"📋 Tables: {len(data.get('tables', []))}\n")
                    
                    f.write("\n")
                
                f.write("=" * 60 + "\n")
                f.write("🎉 Extraction completed successfully!\n")
                f.write("📥 Check your downloads for detailed results.\n")
            
            print("✅ Summary report saved as: extraction_summary.txt")
            
        except Exception as e:
            print(f"❌ Error creating summary report: {str(e)}")

def main():
    """Main execution function"""
    print("🚀 COMPREHENSIVE DOCUMENT EXTRACTION TOOL")
    print("=" * 60)
    print("📋 Supports: PDF, DOCX, PPTX, Excel files")
    print("🔧 Features: Text, Images, Tables, Metadata extraction")
    print("=" * 60)
    
    try:
        extractor = ComprehensiveDocumentExtractor()
        
        # Upload files
        file_paths = extractor.upload_files()
        
        if not file_paths:
            print("❌ No valid files uploaded. Please try again.")
            return
        
        # Process files
        results = extractor.process_files(file_paths)
        
        if results:
            print(f"\n🎉 SUCCESS! Processed {len(results)} file(s)")
            
            # Display quick summary
            print("\n📊 QUICK SUMMARY:")
            print("-" * 40)
            for file_path, data in results.items():
                filename = data.get('filename', os.path.basename(file_path))
                file_type = data.get('file_type', 'Unknown')
                
                if file_type == 'PDF':
                    print(f"📄 {filename}: {data.get('page_count', 0)} pages, {data.get('total_words', 0):,} words, {data.get('total_images', 0)} images")
                elif file_type == 'PPTX':
                    print(f"📄 {filename}: {data.get('total_slides', 0)} slides, {data.get('total_words', 0):,} words")
                elif file_type == 'Excel':
                    print(f"📄 {filename}: {len(data.get('sheets', []))} sheets")
                elif file_type == 'DOCX':
                    print(f"📄 {filename}: {len(data.get('paragraphs', []))} paragraphs, {data.get('total_words', 0):,} words")
            
            # Save results
            print(f"\n💾 Saving extraction results...")
            extractor.save_results(results)
            
            print(f"\n🎉 EXTRACTION COMPLETED SUCCESSFULLY!")
            print("=" * 60)
            print("📥 Downloaded files:")
            print("  • extraction_results.json - Complete extraction data")
            print("  • extraction_summary.txt - Human-readable summary")
            print("📁 Check these folders for additional files:")
            print("  • extracted_images/ - All extracted images")
            print("  • extracted_tables/ - Tables saved as CSV files")
            print("=" * 60)
            
        else:
            print("❌ No files could be processed successfully.")
            print("Please check your files and try again.")
    
    except Exception as e:
        print(f"❌ An error occurred: {str(e)}")
        print("Please try again or check your files.")

# ================================================================
# INSTRUCTIONS FOR GOOGLE COLAB USERS
# ================================================================

print("""
🚀 GOOGLE COLAB DOCUMENT EXTRACTION TOOL
========================================

📋 INSTRUCTIONS:
1. Run the setup cell above first (installs packages)
2. Run the main code cell
3. When prompted, upload your files using the file picker
4. Wait for processing to complete
5. Download the results automatically

📁 SUPPORTED FORMATS:
• PDF files (.pdf)
• Word documents (.docx) 
• PowerPoint presentations (.pptx)
• Excel spreadsheets (.xlsx, .xls)

🔧 FEATURES:
• Complete text extraction with formatting
• Image extraction and saving
• Table extraction using multiple methods
• Metadata extraction
• Font analysis (PDF)
• Comprehensive summary reports

⚡ READY TO START!
Run: main()
""")

# Run the extraction tool
if __name__ == "__main__":
    main()